#!/usr/bin/env python3
"""
Crea slide PowerPoint con loghi reali delle aziende partner
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import math

# Configurazione
SLIDE_WIDTH = Inches(16)
SLIDE_HEIGHT = Inches(9)
OUTPUT_FILE = Path("/Users/bibop/Documents/AI/Lorenz/PowerPoint_Slide/Hyperworks_Partners_Final.pptx")
LOGOS_DIR = Path("/Users/bibop/Documents/AI/Lorenz/PowerPoint_Slide/logos")

# Colori per le categorie
COLORS = {
    "technology": RGBColor(0, 120, 212),      # Blu
    "financial": RGBColor(16, 124, 65),       # Verde
    "operations": RGBColor(232, 89, 12),      # Arancione
    "management": RGBColor(136, 23, 152),     # Viola
}

# Mapping nomi aziende -> file logo
LOGO_FILES = {
    "ThyssenKrupp": "thyssenkrupp.png",
    "Hardt": "hardt.png",
    "Zeleros": "zeleros.png",
    "Swisspod": "swisspod.png",
    "UniCredit": "unicredit.png",
    "Brookfield": "brookfield.png",
    "Macquarie": "macquarie.png",
    "EIB": "eib.png",
    "Italferr": "italferr.png",
    "WeBuild": "webuild.png",
    "Leonardo": "leonardo.png",
    "Fastweb": "fastweb.png",
    "UniTo": "unito.png",
}

CATEGORIES = {
    "technology": {
        "title": "TECHNOLOGY\nSUPPLY CHAIN",
        "companies": ["ThyssenKrupp", "Hardt", "Zeleros", "Swisspod"],
        "color": COLORS["technology"],
    },
    "financial": {
        "title": "FINANCIAL\nSTRUCTURING",
        "companies": ["UniCredit", "Brookfield", "Macquarie", "EIB"],
        "color": COLORS["financial"],
    },
    "operations": {
        "title": "OPERATIONS",
        "companies": ["Italferr"],
        "color": COLORS["operations"],
    },
    "management": {
        "title": "PROJECT\nMANAGEMENT",
        "companies": ["WeBuild", "Leonardo", "Fastweb", "UniTo"],
        "color": COLORS["management"],
    }
}


def add_logo(slide, company_name, x, y, max_width, max_height):
    """Aggiunge un logo mantenendo aspect ratio"""
    logo_file = LOGO_FILES.get(company_name)
    if not logo_file:
        return False

    logo_path = LOGOS_DIR / logo_file
    if not logo_path.exists():
        print(f"  ⚠️  Logo non trovato: {logo_path}")
        return False

    try:
        pic = slide.shapes.add_picture(
            str(logo_path),
            Inches(x), Inches(y),
            width=Inches(max_width)
        )
        # Mantieni aspect ratio entro i limiti
        if pic.height > Inches(max_height):
            ratio = Inches(max_height) / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = Inches(max_height)
        return True
    except Exception as e:
        print(f"  ❌ Errore logo {company_name}: {e}")
        return False


def create_slide_with_logos():
    """Crea slide con loghi reali"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Sfondo
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
    bg.line.fill.background()

    center_x = SLIDE_WIDTH.inches / 2
    center_y = SLIDE_HEIGHT.inches / 2 + 0.2

    # Configurazione cerchi
    circle_configs = [
        (-2.6, -1.4, "technology"),
        (2.6, -1.4, "financial"),
        (-2.6, 1.6, "operations"),
        (2.6, 1.6, "management"),
    ]

    radius = Inches(2.9)

    print("\n📊 Creazione slide con loghi reali...")

    for x_off, y_off, cat_key in circle_configs:
        cat_data = CATEGORIES[cat_key]

        cx = center_x + x_off
        cy = center_y + y_off

        left = Inches(cx) - radius / 2
        top = Inches(cy) - radius / 2

        # Cerchio
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, top, radius, radius
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = cat_data["color"]
        circle.line.color.rgb = cat_data["color"]
        circle.line.width = Pt(2)

        # Titolo categoria
        title_box = slide.shapes.add_textbox(
            left + Inches(0.15),
            top + Inches(0.12),
            radius - Inches(0.3),
            Inches(0.55)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cat_data["title"]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Posiziona loghi aziende
        companies = cat_data["companies"]
        num = len(companies)

        print(f"\n  {cat_key.upper()}:")

        if num == 1:
            # Un solo logo centrato
            logo_w, logo_h = 1.2, 0.5
            add_logo(slide, companies[0],
                    cx - logo_w/2, cy - logo_h/2,
                    logo_w, logo_h)
            print(f"    ✓ {companies[0]}")
        elif num == 4:
            # Griglia 2x2
            logo_w, logo_h = 0.8, 0.4
            spacing_x, spacing_y = 0.9, 0.5
            start_x = cx - spacing_x/2 - logo_w/2
            start_y = cy - 0.3

            for i, company in enumerate(companies):
                row, col = i // 2, i % 2
                lx = start_x + col * spacing_x
                ly = start_y + row * spacing_y

                # Box bianco per logo
                box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(lx - 0.05), Inches(ly - 0.05),
                    Inches(logo_w + 0.1), Inches(logo_h + 0.1)
                )
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                box.line.fill.background()

                if add_logo(slide, company, lx, ly, logo_w, logo_h):
                    print(f"    ✓ {company}")
                else:
                    print(f"    ✗ {company}")

    # Cerchio centrale Hyperworks
    center_radius = Inches(2.1)
    center_left = Inches(center_x) - center_radius / 2
    center_top = Inches(center_y) - center_radius / 2

    center_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_left, center_top,
        center_radius, center_radius
    )
    center_circle.fill.solid()
    center_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    center_circle.line.color.rgb = RGBColor(50, 50, 50)
    center_circle.line.width = Pt(3)

    # Logo Hyperworks
    hw_logo = LOGOS_DIR / "hyperworks.png"
    if hw_logo.exists():
        logo_w = Inches(1.7)
        slide.shapes.add_picture(
            str(hw_logo),
            Inches(center_x) - logo_w / 2,
            Inches(center_y) - Inches(0.2),
            logo_w
        )
        print(f"\n  HYPERWORKS: ✓")

    # Titolo
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        SLIDE_WIDTH - Inches(1), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "HYPERWORKS ECOSYSTEM"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(40, 40, 40)
    p.alignment = PP_ALIGN.CENTER

    # Sottotitolo
    subtitle = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.6),
        SLIDE_WIDTH - Inches(1), Inches(0.35)
    )
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Strategic Partners & Collaborators"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer = slide.shapes.add_textbox(
        Inches(0.5), SLIDE_HEIGHT - Inches(0.5),
        SLIDE_WIDTH - Inches(1), Inches(0.3)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "www.hyper.works"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    prs.save(str(OUTPUT_FILE))
    print(f"\n✅ Slide salvata: {OUTPUT_FILE}")
    return str(OUTPUT_FILE)


if __name__ == "__main__":
    print("=" * 50)
    print("Creazione Slide Hyperworks con Loghi Reali")
    print("=" * 50)
    create_slide_with_logos()
