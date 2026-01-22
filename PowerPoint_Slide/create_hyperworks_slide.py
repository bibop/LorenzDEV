#!/usr/bin/env python3
"""
Crea slide PowerPoint con cerchi sovrapposti per Hyperworks
Design basato sull'immagine di riferimento con 4 aree tematiche
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import math
from pathlib import Path

# In python-pptx 1.0+, RgbColor è ora in pptx.util o si usa direttamente
def RgbColor(r, g, b):
    """Helper per creare colori RGB"""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

# Configurazione
SLIDE_WIDTH = Inches(16)  # Widescreen 16:9
SLIDE_HEIGHT = Inches(9)
OUTPUT_FILE = Path("/Users/bibop/Documents/AI/Lorenz/PowerPoint_Slide/Hyperworks_Partners_Slide.pptx")
LOGOS_DIR = Path("/Users/bibop/Documents/AI/Lorenz/PowerPoint_Slide/logos")

# Colori per le categorie (RGB)
COLORS = {
    "technology": RgbColor(0, 120, 212),      # Blu
    "financial": RgbColor(16, 124, 65),       # Verde
    "operations": RgbColor(232, 89, 12),      # Arancione
    "management": RgbColor(136, 23, 152),     # Viola
    "center": RgbColor(0, 0, 0),              # Nero per Hyperworks
}

# Definizione delle aree e dei loghi
CATEGORIES = {
    "technology": {
        "title": "TECHNOLOGY\nSUPPLY CHAIN",
        "companies": ["ThyssenKrupp", "Hardt", "Zeleros", "Swisspod"],
        "color": COLORS["technology"],
        "position": "top_left"
    },
    "financial": {
        "title": "FINANCIAL\nSTRUCTURING",
        "companies": ["UniCredit", "Brookfield", "Macquarie", "EIB"],
        "color": COLORS["financial"],
        "position": "top_right"
    },
    "operations": {
        "title": "OPERATIONS",
        "companies": ["Italferr"],
        "color": COLORS["operations"],
        "position": "bottom_left"
    },
    "management": {
        "title": "PROJECT\nMANAGEMENT",
        "companies": ["WeBuild", "Leonardo", "Fastweb", "UniTo"],
        "color": COLORS["management"],
        "position": "bottom_right"
    }
}


def create_presentation():
    """Crea la presentazione PowerPoint"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Usa layout blank
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Dimensioni e posizioni dei cerchi
    center_x = SLIDE_WIDTH.inches / 2
    center_y = SLIDE_HEIGHT.inches / 2

    # Cerchi esterni (4 aree)
    outer_radius = Inches(3.0)
    outer_offset = Inches(2.2)  # Distanza dal centro

    # Cerchio centrale (più piccolo)
    center_radius = Inches(1.8)

    # Posizioni dei 4 cerchi esterni
    positions = {
        "top_left": (center_x - outer_offset.inches, center_y - outer_offset.inches * 0.6),
        "top_right": (center_x + outer_offset.inches, center_y - outer_offset.inches * 0.6),
        "bottom_left": (center_x - outer_offset.inches, center_y + outer_offset.inches * 0.6),
        "bottom_right": (center_x + outer_offset.inches, center_y + outer_offset.inches * 0.6),
    }

    # Disegna i cerchi esterni (con trasparenza)
    for cat_key, cat_data in CATEGORIES.items():
        pos = positions[cat_data["position"]]

        # Crea cerchio
        left = Inches(pos[0]) - outer_radius / 2
        top = Inches(pos[1]) - outer_radius / 2

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, top,
            outer_radius, outer_radius
        )

        # Stile del cerchio
        fill = circle.fill
        fill.solid()
        fill.fore_color.rgb = cat_data["color"]

        # Trasparenza (50%)
        circle.fill.fore_color.brightness = 0.3

        # Bordo
        line = circle.line
        line.color.rgb = cat_data["color"]
        line.width = Pt(2)

        # Aggiungi titolo categoria
        title_box = slide.shapes.add_textbox(
            left + Inches(0.3),
            top + Inches(0.3),
            outer_radius - Inches(0.6),
            Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cat_data["title"]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Aggiungi nomi aziende
        companies_y = top + Inches(1.2)
        for i, company in enumerate(cat_data["companies"]):
            comp_box = slide.shapes.add_textbox(
                left + Inches(0.4),
                companies_y + Inches(i * 0.35),
                outer_radius - Inches(0.8),
                Inches(0.35)
            )
            tf = comp_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {company}"
            p.font.size = Pt(11)
            p.font.color.rgb = RgbColor(50, 50, 50)

    # Cerchio centrale (Hyperworks)
    center_left = Inches(center_x) - center_radius / 2
    center_top = Inches(center_y) - center_radius / 2

    center_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_left, center_top,
        center_radius, center_radius
    )

    # Stile cerchio centrale
    center_circle.fill.solid()
    center_circle.fill.fore_color.rgb = RgbColor(255, 255, 255)
    center_circle.line.color.rgb = RgbColor(0, 0, 0)
    center_circle.line.width = Pt(3)

    # Aggiungi logo Hyperworks se esiste
    hyperworks_logo = LOGOS_DIR / "hyperworks.png"
    if hyperworks_logo.exists():
        # Aggiungi immagine al centro
        logo_size = Inches(1.4)
        logo_left = Inches(center_x) - logo_size / 2
        logo_top = Inches(center_y) - logo_size / 2
        slide.shapes.add_picture(
            str(hyperworks_logo),
            logo_left, logo_top,
            logo_size, logo_size
        )
    else:
        # Aggiungi testo placeholder
        text_box = slide.shapes.add_textbox(
            center_left + Inches(0.2),
            center_top + Inches(0.6),
            center_radius - Inches(0.4),
            Inches(0.6)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = "HYPERWORKS"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    # Aggiungi titolo slide
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        SLIDE_WIDTH - Inches(1), Inches(0.6)
    )
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "HYPERWORKS ECOSYSTEM"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RgbColor(30, 30, 30)
    p.alignment = PP_ALIGN.CENTER

    # Aggiungi sottotitolo
    subtitle_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.85),
        SLIDE_WIDTH - Inches(1), Inches(0.4)
    )
    tf = subtitle_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Strategic Partners & Collaborators"
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Salva
    prs.save(str(OUTPUT_FILE))
    print(f"✅ Slide creata: {OUTPUT_FILE}")
    return str(OUTPUT_FILE)


def create_version_with_logo_placeholders():
    """
    Crea una versione alternativa con spazi chiari per i loghi
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Sfondo bianco
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(250, 250, 250)
    background.line.fill.background()

    center_x = SLIDE_WIDTH.inches / 2
    center_y = SLIDE_HEIGHT.inches / 2 + 0.3

    # Layout dei cerchi come nell'immagine originale
    circle_configs = [
        # (x_offset, y_offset, category_key)
        (-2.5, -1.5, "technology"),   # Top-left
        (2.5, -1.5, "financial"),     # Top-right
        (-2.5, 1.5, "operations"),    # Bottom-left
        (2.5, 1.5, "management"),     # Bottom-right
    ]

    radius = Inches(2.8)

    for x_off, y_off, cat_key in circle_configs:
        cat_data = CATEGORIES[cat_key]

        cx = center_x + x_off
        cy = center_y + y_off

        left = Inches(cx) - radius / 2
        top = Inches(cy) - radius / 2

        # Cerchio semi-trasparente
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, top, radius, radius
        )

        circle.fill.solid()
        circle.fill.fore_color.rgb = cat_data["color"]
        circle.line.color.rgb = cat_data["color"]
        circle.line.width = Pt(2)

        # Titolo categoria (in alto nel cerchio)
        title_box = slide.shapes.add_textbox(
            left + Inches(0.2),
            top + Inches(0.15),
            radius - Inches(0.4),
            Inches(0.6)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cat_data["title"]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Griglia per loghi delle aziende
        num_companies = len(cat_data["companies"])
        logo_size = Inches(0.5)

        if num_companies <= 2:
            cols = num_companies
        else:
            cols = 2
        rows = math.ceil(num_companies / cols)

        start_x = cx - (cols * 0.6) / 2
        start_y = cy - 0.1

        for i, company in enumerate(cat_data["companies"]):
            row = i // cols
            col = i % cols

            lx = Inches(start_x + col * 0.65)
            ly = Inches(start_y + row * 0.55)

            # Box placeholder per logo
            logo_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                lx, ly,
                Inches(0.55), Inches(0.45)
            )
            logo_box.fill.solid()
            logo_box.fill.fore_color.rgb = RgbColor(255, 255, 255)
            logo_box.line.color.rgb = RgbColor(200, 200, 200)
            logo_box.line.width = Pt(1)

            # Nome azienda sotto il box
            name_box = slide.shapes.add_textbox(
                lx - Inches(0.1), ly + Inches(0.45),
                Inches(0.75), Inches(0.25)
            )
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = company
            p.font.size = Pt(8)
            p.font.color.rgb = RgbColor(60, 60, 60)
            p.alignment = PP_ALIGN.CENTER

    # Cerchio centrale Hyperworks (più grande, bianco)
    center_radius = Inches(2.0)
    center_left = Inches(center_x) - center_radius / 2
    center_top = Inches(center_y) - center_radius / 2

    center_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_left, center_top,
        center_radius, center_radius
    )
    center_circle.fill.solid()
    center_circle.fill.fore_color.rgb = RgbColor(255, 255, 255)
    center_circle.line.color.rgb = RgbColor(30, 30, 30)
    center_circle.line.width = Pt(3)

    # Logo Hyperworks al centro
    hyperworks_logo = LOGOS_DIR / "hyperworks.png"
    if hyperworks_logo.exists():
        logo_w = Inches(1.6)
        logo_h = Inches(0.36)  # Aspect ratio ~4.5:1
        slide.shapes.add_picture(
            str(hyperworks_logo),
            Inches(center_x) - logo_w / 2,
            Inches(center_y) - logo_h / 2,
            logo_w
        )
    else:
        text_box = slide.shapes.add_textbox(
            center_left + Inches(0.3),
            center_top + Inches(0.8),
            center_radius - Inches(0.6),
            Inches(0.5)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = "HYPERWORKS"
        p.font.bold = True
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER

    # Titolo
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        SLIDE_WIDTH - Inches(1), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "HYPERWORKS ECOSYSTEM"
    p.font.bold = True
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER

    output_v2 = OUTPUT_FILE.parent / "Hyperworks_Partners_v2.pptx"
    prs.save(str(output_v2))
    print(f"✅ Slide v2 creata: {output_v2}")
    return str(output_v2)


if __name__ == "__main__":
    print("=" * 60)
    print("Creazione Slide Hyperworks Partners")
    print("=" * 60)

    # Versione 1: Cerchi sovrapposti base
    create_presentation()

    # Versione 2: Con placeholder per loghi
    create_version_with_logo_placeholders()

    print("\n📝 Note:")
    print("   - I loghi delle aziende devono essere aggiunti manualmente")
    print("   - Usa 'Inserisci > Immagine' in PowerPoint")
    print("   - Assicurati che i loghi siano PNG con sfondo trasparente")
    print("\n🔗 Fonti per scaricare loghi HD:")
    print("   - https://seeklogo.com")
    print("   - https://worldvectorlogo.com")
    print("   - https://brandsoftheworld.com")
    print("   - Siti ufficiali aziende (sezione Press/Media)")
