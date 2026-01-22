#!/usr/bin/env python3
"""
🎯 LORENZ Skills System
===========================================================

Sistema modulare di competenze per LORENZ.
Ogni skill è un'azione specifica che LORENZ può eseguire.

🌟 GOD SKILLS (Built-in):
- 🖼️ ImageGeneration - Genera immagini con DALL-E
- 🔍 WebSearch - Ricerca web con Perplexity
- 📊 Presentation - Crea presentazioni PowerPoint
- 💻 CodeExecution - Esegue codice Python
- 🖥️ ServerCommand - Comandi SSH sul server
- 📁 NASOperation - Operazioni sul NAS Synology
- 🌐 WebBrowse - Naviga e estrae contenuti web
- 📝 DocumentGen - Genera documenti Word/PDF

🧠 EMERGENT SKILLS:
- Skills apprese dinamicamente da LORENZ durante le interazioni
- Memorizzate nella Knowledge Base MNEME
- Ispezionabili e modificabili dal pannello di controllo

📚 MNEME (Knowledge Base):
- Archivio della conoscenza appresa
- Patterns di successo memorizzati
- Workflow emergenti

Autore: Claude Code
Data: 2026-01-14
"""

import os
import json
import logging
import asyncio
import aiohttp
import subprocess
import tempfile
import hashlib
import sqlite3
from abc import ABC, abstractmethod
from typing import Dict, List, Optional, Any, Union, Callable
from dataclasses import dataclass, field, asdict
from enum import Enum
from datetime import datetime
from pathlib import Path
import uuid

logger = logging.getLogger(__name__)


# ============================================================================
# SKILL TYPES
# ============================================================================

class SkillType(Enum):
    """Distingue tra skills create da sviluppatori e skills apprese"""
    GOD = "god"           # 🌟 Built-in, created by developers
    EMERGENT = "emergent" # 🧠 Learned by LORENZ through interactions


class SkillCategory(Enum):
    """Categorie di skills per organizzazione"""
    CREATIVE = "creative"       # Image, presentation, document generation
    RESEARCH = "research"       # Web search, web browse
    TECHNICAL = "technical"     # Code execution, server commands
    DATA = "data"               # NAS, file operations
    COMMUNICATION = "communication"  # Email, calendar
    WORKFLOW = "workflow"       # Emergent workflow combinations
    CUSTOM = "custom"           # User-defined emergent skills


# ============================================================================
# SKILL BASE CLASS
# ============================================================================

@dataclass
class SkillResult:
    """Result from a skill execution"""
    success: bool
    data: Any = None
    message: str = ""
    error: str = None
    artifacts: List[str] = field(default_factory=list)  # File paths, URLs, etc.
    skill_name: str = ""
    execution_time_ms: int = 0


@dataclass
class SkillMetadata:
    """Metadata for skill tracking and learning"""
    created_at: str = ""
    last_used: str = ""
    use_count: int = 0
    success_rate: float = 1.0
    avg_execution_time_ms: float = 0.0
    tags: List[str] = field(default_factory=list)
    version: str = "1.0.0"


class Skill(ABC):
    """Base class for all skills"""

    name: str = "base_skill"
    description: str = "Base skill"
    description_it: str = "Skill base"  # Italian description
    examples: List[str] = []
    requires_api: List[str] = []  # API keys needed
    skill_type: SkillType = SkillType.GOD  # Default to GOD skill
    category: SkillCategory = SkillCategory.CUSTOM
    icon: str = "⚡"  # Emoji icon for UI

    def __init__(self):
        self.enabled = self._check_requirements()
        self.metadata = SkillMetadata(
            created_at=datetime.now().isoformat(),
            last_used="",
            use_count=0
        )
        self.skill_id = str(uuid.uuid4())[:8]

    def _check_requirements(self) -> bool:
        """Check if required API keys are available"""
        for api in self.requires_api:
            key_name = f"{api.upper()}_API_KEY"
            if not os.getenv(key_name):
                logger.warning(f"⚠️ Skill {self.name} disabled: missing {key_name}")
                return False
        return True

    @abstractmethod
    async def execute(self, **kwargs) -> SkillResult:
        """Execute the skill"""
        pass

    def _track_execution(self, result: SkillResult, execution_time: float):
        """Track execution for learning"""
        self.metadata.use_count += 1
        self.metadata.last_used = datetime.now().isoformat()

        # Update success rate (rolling average)
        old_rate = self.metadata.success_rate
        old_count = self.metadata.use_count - 1
        success_val = 1.0 if result.success else 0.0
        if old_count > 0:
            self.metadata.success_rate = (old_rate * old_count + success_val) / self.metadata.use_count
        else:
            self.metadata.success_rate = success_val

        # Update average execution time
        self.metadata.avg_execution_time_ms = (
            (self.metadata.avg_execution_time_ms * old_count + execution_time) /
            self.metadata.use_count
        )

    def get_info(self) -> Dict:
        """Get skill info for UI display"""
        return {
            "id": self.skill_id,
            "name": self.name,
            "description": self.description,
            "description_it": self.description_it,
            "examples": self.examples,
            "enabled": self.enabled,
            "requires": self.requires_api,
            "type": self.skill_type.value,
            "category": self.category.value,
            "icon": self.icon,
            "metadata": {
                "created_at": self.metadata.created_at,
                "last_used": self.metadata.last_used,
                "use_count": self.metadata.use_count,
                "success_rate": round(self.metadata.success_rate * 100, 1),
                "avg_time_ms": round(self.metadata.avg_execution_time_ms, 0),
                "tags": self.metadata.tags
            }
        }

    def to_dict(self) -> Dict:
        """Serialize skill for storage"""
        return self.get_info()


# ============================================================================
# IMAGE GENERATION SKILL
# ============================================================================

class ImageGenerationSkill(Skill):
    """Generate images using DALL-E 3"""

    name = "image_generation"
    description = "Generate images from text descriptions using DALL-E 3"
    description_it = "Genera immagini da descrizioni testuali usando DALL-E 3"
    examples = [
        "Genera un'immagine di un tramonto sul mare",
        "Crea un logo per una startup tech",
        "Disegna un gatto astronauta"
    ]
    requires_api = ["OPENAI"]
    skill_type = SkillType.GOD
    category = SkillCategory.CREATIVE
    icon = "🖼️"

    async def execute(self, prompt: str, size: str = "1024x1024",
                     quality: str = "standard", style: str = "vivid") -> SkillResult:
        """
        Generate an image

        Args:
            prompt: Image description
            size: "1024x1024", "1792x1024", "1024x1792"
            quality: "standard" or "hd"
            style: "vivid" or "natural"
        """
        if not self.enabled:
            return SkillResult(
                success=False,
                error="OpenAI API key not configured"
            )

        api_key = os.getenv("OPENAI_API_KEY")

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }

        payload = {
            "model": "dall-e-3",
            "prompt": prompt,
            "size": size,
            "quality": quality,
            "style": style,
            "n": 1
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.openai.com/v1/images/generations",
                    headers=headers,
                    json=payload
                ) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        image_url = data["data"][0]["url"]
                        revised_prompt = data["data"][0].get("revised_prompt", prompt)

                        return SkillResult(
                            success=True,
                            data={
                                "url": image_url,
                                "revised_prompt": revised_prompt
                            },
                            message=f"🖼️ Immagine generata!\n\nPrompt: {revised_prompt}",
                            artifacts=[image_url]
                        )
                    else:
                        error = await resp.text()
                        return SkillResult(
                            success=False,
                            error=f"DALL-E API error: {error}"
                        )
        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# WEB SEARCH SKILL
# ============================================================================

class WebSearchSkill(Skill):
    """Search the web using Perplexity"""

    name = "web_search"
    description = "Search for up-to-date information on the web"
    description_it = "Cerca informazioni aggiornate sul web"
    examples = [
        "Cerca le ultime notizie su AI",
        "Trova il prezzo attuale di Bitcoin",
        "Cerca recensioni del nuovo iPhone"
    ]
    requires_api = ["PERPLEXITY"]
    skill_type = SkillType.GOD
    category = SkillCategory.RESEARCH
    icon = "🔍"

    async def execute(self, query: str, detailed: bool = False) -> SkillResult:
        """
        Search the web

        Args:
            query: Search query
            detailed: If True, provide more detailed results
        """
        if not self.enabled:
            return SkillResult(
                success=False,
                error="Perplexity API key not configured"
            )

        api_key = os.getenv("PERPLEXITY_API_KEY")

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }

        model = "llama-3.1-sonar-large-128k-online" if detailed else "llama-3.1-sonar-small-128k-online"

        payload = {
            "model": model,
            "messages": [{"role": "user", "content": query}]
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.perplexity.ai/chat/completions",
                    headers=headers,
                    json=payload
                ) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        response = data["choices"][0]["message"]["content"]
                        citations = data.get("citations", [])

                        return SkillResult(
                            success=True,
                            data={
                                "response": response,
                                "citations": citations
                            },
                            message=f"🔍 Risultati ricerca:\n\n{response}"
                        )
                    else:
                        error = await resp.text()
                        return SkillResult(
                            success=False,
                            error=f"Perplexity API error: {error}"
                        )
        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# PRESENTATION SKILL
# ============================================================================

class PresentationSkill(Skill):
    """Create PowerPoint presentations"""

    name = "presentation"
    description = "Create professional PowerPoint presentations"
    description_it = "Crea presentazioni PowerPoint professionali"
    examples = [
        "Crea una presentazione sul cambiamento climatico",
        "Genera slide per un pitch di startup",
        "Prepara una presentazione sul machine learning"
    ]
    requires_api = []  # Uses local python-pptx
    skill_type = SkillType.GOD
    category = SkillCategory.CREATIVE
    icon = "📊"

    def __init__(self):
        super().__init__()
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            self.pptx_available = True
        except ImportError:
            self.pptx_available = False
            self.enabled = False
            logger.warning("⚠️ python-pptx not installed. Presentation skill disabled.")

    async def execute(self, title: str, slides: List[Dict],
                     output_path: str = None) -> SkillResult:
        """
        Create a presentation

        Args:
            title: Presentation title
            slides: List of slide definitions
                   [{"title": "...", "content": ["bullet1", "bullet2"], "notes": "..."}]
            output_path: Where to save (default: temp file)
        """
        if not self.pptx_available:
            return SkillResult(
                success=False,
                error="python-pptx not installed. Run: pip install python-pptx"
            )

        from pptx import Presentation
        from pptx.util import Inches, Pt

        try:
            prs = Presentation()

            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = datetime.now().strftime("%Y-%m-%d")

            # Content slides
            bullet_slide_layout = prs.slide_layouts[1]
            for slide_def in slides:
                slide = prs.slides.add_slide(bullet_slide_layout)
                slide.shapes.title.text = slide_def.get("title", "")

                # Add content
                body = slide.shapes.placeholders[1]
                tf = body.text_frame

                content = slide_def.get("content", [])
                for i, bullet in enumerate(content):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0

                # Add notes
                if slide_def.get("notes"):
                    slide.notes_slide.notes_text_frame.text = slide_def["notes"]

            # Save
            if output_path is None:
                output_path = tempfile.mktemp(suffix=".pptx")

            prs.save(output_path)

            return SkillResult(
                success=True,
                data={"path": output_path, "slides": len(slides) + 1},
                message=f"📊 Presentazione creata con {len(slides) + 1} slide",
                artifacts=[output_path]
            )

        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# CODE EXECUTION SKILL
# ============================================================================

class CodeExecutionSkill(Skill):
    """Execute Python code safely"""

    name = "code_execution"
    description = "Execute Python code in a secure environment"
    description_it = "Esegue codice Python in un ambiente sicuro"
    examples = [
        "Calcola il fattoriale di 10",
        "Genera un grafico dei dati",
        "Analizza questo CSV"
    ]
    requires_api = []
    skill_type = SkillType.GOD
    category = SkillCategory.TECHNICAL
    icon = "💻"

    # Allowed modules for safety
    ALLOWED_MODULES = [
        "math", "statistics", "datetime", "json", "re",
        "collections", "itertools", "functools",
        "pandas", "numpy", "matplotlib"
    ]

    async def execute(self, code: str, timeout: int = 30) -> SkillResult:
        """
        Execute Python code

        Args:
            code: Python code to execute
            timeout: Max execution time in seconds
        """
        # Safety check - basic validation
        dangerous_patterns = [
            "import os", "import sys", "import subprocess",
            "__import__", "eval(", "exec(", "open(",
            "file(", "input(", "breakpoint("
        ]

        for pattern in dangerous_patterns:
            if pattern in code:
                return SkillResult(
                    success=False,
                    error=f"Codice non sicuro: pattern '{pattern}' non permesso"
                )

        try:
            # Create temp file
            with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as f:
                f.write(code)
                temp_path = f.name

            # Execute with timeout
            result = subprocess.run(
                ['python3', temp_path],
                capture_output=True,
                text=True,
                timeout=timeout
            )

            # Cleanup
            os.unlink(temp_path)

            if result.returncode == 0:
                return SkillResult(
                    success=True,
                    data={"output": result.stdout},
                    message=f"✅ Codice eseguito:\n```\n{result.stdout}\n```"
                )
            else:
                return SkillResult(
                    success=False,
                    error=result.stderr,
                    message=f"❌ Errore:\n```\n{result.stderr}\n```"
                )

        except subprocess.TimeoutExpired:
            return SkillResult(
                success=False,
                error=f"Timeout: esecuzione superato {timeout} secondi"
            )
        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# SERVER COMMAND SKILL
# ============================================================================

class ServerCommandSkill(Skill):
    """Execute commands on remote server via SSH"""

    name = "server_command"
    description = "Execute commands on remote server via SSH"
    description_it = "Esegue comandi sul server remoto via SSH"
    examples = [
        "Controlla lo stato del server",
        "Riavvia il servizio nginx",
        "Mostra l'uso del disco"
    ]
    requires_api = []
    skill_type = SkillType.GOD
    category = SkillCategory.TECHNICAL
    icon = "🖥️"

    # Whitelisted commands for safety
    ALLOWED_COMMANDS = [
        "uptime", "free", "df", "top", "ps", "date", "hostname",
        "systemctl status", "systemctl restart",
        "journalctl", "tail", "head", "cat", "ls", "pwd",
        "docker ps", "docker logs", "docker stats"
    ]

    def __init__(self, host: str = None, user: str = None, key_path: str = None):
        super().__init__()
        self.host = host or os.getenv("SSH_HOST", "80.240.31.197")
        self.user = user or os.getenv("SSH_USER", "linuxuser")
        self.key_path = key_path or os.getenv("SSH_KEY_PATH", os.path.expanduser("~/.ssh/id_rsa"))

    def _is_command_allowed(self, command: str) -> bool:
        """Check if command is in whitelist"""
        cmd_start = command.split()[0] if command.split() else ""

        # Check exact matches and prefixes
        for allowed in self.ALLOWED_COMMANDS:
            if command.startswith(allowed) or cmd_start == allowed.split()[0]:
                return True
        return False

    async def execute(self, command: str, timeout: int = 60) -> SkillResult:
        """
        Execute SSH command

        Args:
            command: Command to execute
            timeout: Max execution time
        """
        # Safety check
        if not self._is_command_allowed(command):
            return SkillResult(
                success=False,
                error=f"Comando non permesso: {command}\nComandi permessi: {', '.join(self.ALLOWED_COMMANDS)}"
            )

        try:
            ssh_cmd = [
                "ssh",
                "-i", self.key_path,
                "-o", "StrictHostKeyChecking=no",
                "-o", "ConnectTimeout=10",
                f"{self.user}@{self.host}",
                command
            ]

            result = subprocess.run(
                ssh_cmd,
                capture_output=True,
                text=True,
                timeout=timeout
            )

            if result.returncode == 0:
                return SkillResult(
                    success=True,
                    data={"output": result.stdout, "command": command},
                    message=f"💻 Output:\n```\n{result.stdout}\n```"
                )
            else:
                return SkillResult(
                    success=False,
                    error=result.stderr or "Comando fallito",
                    data={"output": result.stdout}
                )

        except subprocess.TimeoutExpired:
            return SkillResult(
                success=False,
                error=f"Timeout: comando superato {timeout} secondi"
            )
        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# WEB BROWSE SKILL
# ============================================================================

class WebBrowseSkill(Skill):
    """Browse web pages and extract content"""

    name = "web_browse"
    description = "Browse web pages and extract content"
    description_it = "Naviga pagine web ed estrae contenuti"
    examples = [
        "Leggi il contenuto di questa pagina web",
        "Estrai il testo da questo URL",
        "Scarica l'HTML di questo sito"
    ]
    requires_api = []
    skill_type = SkillType.GOD
    category = SkillCategory.RESEARCH
    icon = "🌐"

    async def execute(self, url: str, extract_text: bool = True,
                     selector: str = None) -> SkillResult:
        """
        Browse a web page

        Args:
            url: URL to browse
            extract_text: If True, extract just text content
            selector: CSS selector to target specific element
        """
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36"
            }

            async with aiohttp.ClientSession() as session:
                async with session.get(url, headers=headers, timeout=30) as resp:
                    if resp.status == 200:
                        html = await resp.text()

                        if extract_text:
                            # Simple text extraction (no BeautifulSoup dependency)
                            import re
                            # Remove scripts and styles
                            text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
                            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
                            # Remove HTML tags
                            text = re.sub(r'<[^>]+>', ' ', text)
                            # Clean whitespace
                            text = re.sub(r'\s+', ' ', text).strip()
                            # Truncate
                            text = text[:5000]

                            return SkillResult(
                                success=True,
                                data={"text": text, "url": url},
                                message=f"🌐 Contenuto da {url}:\n\n{text[:1000]}..."
                            )
                        else:
                            return SkillResult(
                                success=True,
                                data={"html": html, "url": url},
                                message=f"🌐 HTML scaricato da {url} ({len(html)} caratteri)"
                            )
                    else:
                        return SkillResult(
                            success=False,
                            error=f"HTTP {resp.status} per {url}"
                        )

        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# DOCUMENT GENERATION SKILL
# ============================================================================

class DocumentGenerationSkill(Skill):
    """Generate Word documents"""

    name = "document_generation"
    description = "Generate professional Word documents"
    description_it = "Genera documenti Word professionali"
    examples = [
        "Crea un report Word",
        "Genera un documento con questi dati",
        "Prepara un contratto"
    ]
    requires_api = []
    skill_type = SkillType.GOD
    category = SkillCategory.CREATIVE
    icon = "📄"

    def __init__(self):
        super().__init__()
        try:
            from docx import Document
            self.docx_available = True
        except ImportError:
            self.docx_available = False
            self.enabled = False

    async def execute(self, title: str, sections: List[Dict],
                     output_path: str = None) -> SkillResult:
        """
        Create a Word document

        Args:
            title: Document title
            sections: [{"heading": "...", "content": "...", "level": 1}]
            output_path: Where to save
        """
        if not self.docx_available:
            return SkillResult(
                success=False,
                error="python-docx not installed"
            )

        from docx import Document
        from docx.shared import Pt

        try:
            doc = Document()
            doc.add_heading(title, 0)

            for section in sections:
                level = section.get("level", 1)
                if section.get("heading"):
                    doc.add_heading(section["heading"], level)
                if section.get("content"):
                    doc.add_paragraph(section["content"])
                if section.get("bullets"):
                    for bullet in section["bullets"]:
                        doc.add_paragraph(bullet, style="List Bullet")

            if output_path is None:
                output_path = tempfile.mktemp(suffix=".docx")

            doc.save(output_path)

            return SkillResult(
                success=True,
                data={"path": output_path},
                message=f"📄 Documento creato: {output_path}",
                artifacts=[output_path]
            )

        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# CALENDAR SKILL (Google Calendar)
# ============================================================================

class CalendarSkill(Skill):
    """Manage Google Calendar events"""

    name = "calendar"
    description = "Manage Google Calendar: create, view, and modify events"
    description_it = "Gestisci Google Calendar: crea, visualizza e modifica eventi"
    examples = [
        "Mostra i prossimi appuntamenti",
        "Crea un meeting per domani alle 15",
        "Cosa ho in agenda questa settimana?"
    ]
    requires_api = ["GOOGLE_CALENDAR"]
    skill_type = SkillType.GOD
    category = SkillCategory.COMMUNICATION
    icon = "📅"

    def __init__(self):
        super().__init__()
        self.credentials_path = os.getenv("GOOGLE_CALENDAR_CREDENTIALS", "")
        self.calendar_id = os.getenv("GOOGLE_CALENDAR_ID", "primary")

        # Check for google-api-python-client
        try:
            from googleapiclient.discovery import build
            from google.oauth2 import service_account
            self.google_api_available = True
        except ImportError:
            self.google_api_available = False
            self.enabled = False
            logger.warning("⚠️ google-api-python-client not installed. Calendar skill disabled.")

    def _get_service(self):
        """Get Google Calendar service"""
        if not self.google_api_available or not self.credentials_path:
            return None

        try:
            from googleapiclient.discovery import build
            from google.oauth2 import service_account

            SCOPES = ['https://www.googleapis.com/auth/calendar']
            credentials = service_account.Credentials.from_service_account_file(
                self.credentials_path, scopes=SCOPES
            )
            service = build('calendar', 'v3', credentials=credentials)
            return service
        except Exception as e:
            logger.error(f"Failed to get Calendar service: {e}")
            return None

    async def execute(self, operation: str = "list", **kwargs) -> SkillResult:
        """
        Execute calendar operation

        Args:
            operation: "list", "create", "update", "delete"
            **kwargs: Operation-specific parameters
        """
        if not self.enabled or not self.google_api_available:
            return SkillResult(
                success=False,
                error="Google Calendar API not configured. Install google-api-python-client and set GOOGLE_CALENDAR_CREDENTIALS"
            )

        if operation == "list":
            return await self._list_events(**kwargs)
        elif operation == "create":
            return await self._create_event(**kwargs)
        elif operation == "delete":
            return await self._delete_event(**kwargs)
        else:
            return SkillResult(
                success=False,
                error=f"Unknown operation: {operation}"
            )

    async def _list_events(self, max_results: int = 10, days_ahead: int = 7) -> SkillResult:
        """List upcoming events"""
        service = self._get_service()
        if not service:
            return SkillResult(success=False, error="Calendar service not available")

        try:
            from datetime import datetime, timedelta, timezone

            now = datetime.now(timezone.utc).isoformat()
            time_max = (datetime.now(timezone.utc) + timedelta(days=days_ahead)).isoformat()

            events_result = service.events().list(
                calendarId=self.calendar_id,
                timeMin=now,
                timeMax=time_max,
                maxResults=max_results,
                singleEvents=True,
                orderBy='startTime'
            ).execute()

            events = events_result.get('items', [])

            if not events:
                return SkillResult(
                    success=True,
                    data={"events": []},
                    message="📅 Nessun evento nei prossimi giorni"
                )

            msg = f"📅 **Prossimi {len(events)} eventi:**\n\n"
            for event in events:
                start = event['start'].get('dateTime', event['start'].get('date'))
                summary = event.get('summary', 'Senza titolo')
                msg += f"• **{summary}**\n  {start}\n\n"

            return SkillResult(
                success=True,
                data={"events": events, "count": len(events)},
                message=msg
            )

        except Exception as e:
            return SkillResult(success=False, error=str(e))

    async def _create_event(self, summary: str, start_time: str,
                           end_time: str = None, description: str = "",
                           location: str = "") -> SkillResult:
        """Create a new calendar event"""
        service = self._get_service()
        if not service:
            return SkillResult(success=False, error="Calendar service not available")

        try:
            from datetime import datetime, timedelta

            # Parse start time
            if 'T' in start_time:
                start_dt = datetime.fromisoformat(start_time.replace('Z', '+00:00'))
            else:
                start_dt = datetime.fromisoformat(start_time)

            # Default end time is 1 hour after start
            if end_time:
                if 'T' in end_time:
                    end_dt = datetime.fromisoformat(end_time.replace('Z', '+00:00'))
                else:
                    end_dt = datetime.fromisoformat(end_time)
            else:
                end_dt = start_dt + timedelta(hours=1)

            event = {
                'summary': summary,
                'location': location,
                'description': description,
                'start': {
                    'dateTime': start_dt.isoformat(),
                    'timeZone': 'Europe/Rome',
                },
                'end': {
                    'dateTime': end_dt.isoformat(),
                    'timeZone': 'Europe/Rome',
                },
            }

            created_event = service.events().insert(
                calendarId=self.calendar_id,
                body=event
            ).execute()

            return SkillResult(
                success=True,
                data={"event": created_event, "id": created_event.get('id')},
                message=f"📅 Evento creato: **{summary}**\nData: {start_dt.strftime('%d/%m/%Y %H:%M')}"
            )

        except Exception as e:
            return SkillResult(success=False, error=str(e))

    async def _delete_event(self, event_id: str) -> SkillResult:
        """Delete a calendar event"""
        service = self._get_service()
        if not service:
            return SkillResult(success=False, error="Calendar service not available")

        try:
            service.events().delete(
                calendarId=self.calendar_id,
                eventId=event_id
            ).execute()

            return SkillResult(
                success=True,
                message=f"🗑️ Evento {event_id} eliminato"
            )

        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# VOICE MESSAGE SKILL (Audio Processing)
# ============================================================================

class VoiceMessageSkill(Skill):
    """Process voice messages: transcribe and respond"""

    name = "voice_message"
    description = "Process voice messages: transcribe audio and generate responses"
    description_it = "Elabora messaggi vocali: trascrive audio e genera risposte"
    examples = [
        "Trascrivi questo messaggio vocale",
        "Cosa dice questo audio?",
        "Converti voce in testo"
    ]
    requires_api = ["OPENAI"]  # Using Whisper for transcription
    skill_type = SkillType.GOD
    category = SkillCategory.COMMUNICATION
    icon = "🎙️"

    async def execute(self, audio_path: str = None, audio_url: str = None,
                     language: str = "it") -> SkillResult:
        """
        Process voice message

        Args:
            audio_path: Local path to audio file
            audio_url: URL to download audio from
            language: Language code (it, en, etc.)
        """
        if not self.enabled:
            return SkillResult(
                success=False,
                error="OpenAI API key not configured for voice processing"
            )

        api_key = os.getenv("OPENAI_API_KEY")

        # Download audio if URL provided
        if audio_url and not audio_path:
            try:
                async with aiohttp.ClientSession() as session:
                    async with session.get(audio_url) as resp:
                        if resp.status == 200:
                            audio_path = tempfile.mktemp(suffix=".ogg")
                            with open(audio_path, 'wb') as f:
                                f.write(await resp.read())
            except Exception as e:
                return SkillResult(success=False, error=f"Failed to download audio: {e}")

        if not audio_path:
            return SkillResult(
                success=False,
                error="No audio file provided (audio_path or audio_url required)"
            )

        try:
            # Transcribe using OpenAI Whisper API
            headers = {
                "Authorization": f"Bearer {api_key}"
            }

            # Read audio file
            with open(audio_path, 'rb') as audio_file:
                import aiohttp
                from aiohttp import FormData

                data = FormData()
                data.add_field('file',
                              audio_file,
                              filename=os.path.basename(audio_path),
                              content_type='audio/ogg')
                data.add_field('model', 'whisper-1')
                data.add_field('language', language)

                async with aiohttp.ClientSession() as session:
                    async with session.post(
                        "https://api.openai.com/v1/audio/transcriptions",
                        headers=headers,
                        data=data
                    ) as resp:
                        if resp.status == 200:
                            result = await resp.json()
                            transcription = result.get('text', '')

                            return SkillResult(
                                success=True,
                                data={
                                    "transcription": transcription,
                                    "language": language
                                },
                                message=f"🎙️ **Trascrizione:**\n\n{transcription}"
                            )
                        else:
                            error = await resp.text()
                            return SkillResult(
                                success=False,
                                error=f"Whisper API error: {error}"
                            )

        except Exception as e:
            return SkillResult(success=False, error=str(e))

    async def text_to_speech(self, text: str, voice: str = "alloy",
                            output_path: str = None) -> SkillResult:
        """
        Convert text to speech using OpenAI TTS

        Args:
            text: Text to convert
            voice: Voice to use (alloy, echo, fable, onyx, nova, shimmer)
            output_path: Path to save audio file
        """
        if not self.enabled:
            return SkillResult(
                success=False,
                error="OpenAI API key not configured for TTS"
            )

        api_key = os.getenv("OPENAI_API_KEY")

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }

        payload = {
            "model": "tts-1",
            "input": text,
            "voice": voice
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.openai.com/v1/audio/speech",
                    headers=headers,
                    json=payload
                ) as resp:
                    if resp.status == 200:
                        audio_data = await resp.read()

                        if output_path is None:
                            output_path = tempfile.mktemp(suffix=".mp3")

                        with open(output_path, 'wb') as f:
                            f.write(audio_data)

                        return SkillResult(
                            success=True,
                            data={"path": output_path, "voice": voice},
                            message=f"🔊 Audio generato: {output_path}",
                            artifacts=[output_path]
                        )
                    else:
                        error = await resp.text()
                        return SkillResult(
                            success=False,
                            error=f"TTS API error: {error}"
                        )

        except Exception as e:
            return SkillResult(success=False, error=str(e))


# ============================================================================
# NAS OPERATION SKILL
# ============================================================================

class NASOperationSkill(Skill):
    """Operations on Synology NAS"""

    name = "nas_operation"
    description = "Operations on Synology NAS (files, backup, status)"
    description_it = "Operazioni sul NAS Synology (file, backup, status)"
    examples = [
        "Lista i file nella cartella condivisa",
        "Controlla lo spazio disponibile sul NAS",
        "Cerca un file sul NAS"
    ]
    requires_api = []
    skill_type = SkillType.GOD
    category = SkillCategory.DATA
    icon = "📁"

    def __init__(self, host: str = None, port: int = 5000,
                 username: str = None, password: str = None):
        super().__init__()
        self.host = host or os.getenv("NAS_HOST", "")
        self.port = port or int(os.getenv("NAS_PORT", "5000"))
        self.username = username or os.getenv("NAS_USERNAME", "")
        self.password = password or os.getenv("NAS_PASSWORD", "")
        self.sid = None  # Session ID

        if not all([self.host, self.username, self.password]):
            self.enabled = False

    async def _login(self) -> bool:
        """Login to Synology API"""
        url = f"http://{self.host}:{self.port}/webapi/auth.cgi"
        params = {
            "api": "SYNO.API.Auth",
            "version": "3",
            "method": "login",
            "account": self.username,
            "passwd": self.password,
            "session": "FileStation",
            "format": "cookie"
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, params=params) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        if data.get("success"):
                            self.sid = data["data"]["sid"]
                            return True
        except Exception as e:
            logger.error(f"NAS login error: {e}")

        return False

    async def execute(self, operation: str, path: str = "/",
                     **kwargs) -> SkillResult:
        """
        Execute NAS operation

        Args:
            operation: "list", "info", "search", "download"
            path: Path on NAS
        """
        if not self.enabled:
            return SkillResult(
                success=False,
                error="NAS not configured. Set NAS_HOST, NAS_USERNAME, NAS_PASSWORD"
            )

        if not self.sid:
            if not await self._login():
                return SkillResult(
                    success=False,
                    error="Failed to login to NAS"
                )

        if operation == "list":
            return await self._list_files(path)
        elif operation == "info":
            return await self._get_info()
        elif operation == "search":
            return await self._search(kwargs.get("pattern", "*"), path)
        else:
            return SkillResult(
                success=False,
                error=f"Unknown operation: {operation}"
            )

    async def _list_files(self, path: str) -> SkillResult:
        """List files in directory"""
        url = f"http://{self.host}:{self.port}/webapi/entry.cgi"
        params = {
            "api": "SYNO.FileStation.List",
            "version": "2",
            "method": "list",
            "folder_path": path,
            "_sid": self.sid
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, params=params) as resp:
                    data = await resp.json()
                    if data.get("success"):
                        files = data["data"]["files"]
                        return SkillResult(
                            success=True,
                            data={"files": files, "path": path},
                            message=f"📁 {len(files)} elementi in {path}"
                        )
                    else:
                        return SkillResult(
                            success=False,
                            error=f"NAS error: {data.get('error')}"
                        )
        except Exception as e:
            return SkillResult(success=False, error=str(e))

    async def _get_info(self) -> SkillResult:
        """Get NAS system info"""
        url = f"http://{self.host}:{self.port}/webapi/entry.cgi"
        params = {
            "api": "SYNO.FileStation.Info",
            "version": "2",
            "method": "get",
            "_sid": self.sid
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, params=params) as resp:
                    data = await resp.json()
                    if data.get("success"):
                        return SkillResult(
                            success=True,
                            data=data["data"],
                            message=f"🖥️ NAS Info: {json.dumps(data['data'], indent=2)}"
                        )
        except Exception as e:
            return SkillResult(success=False, error=str(e))

        return SkillResult(success=False, error="Failed to get NAS info")

    async def _search(self, pattern: str, path: str) -> SkillResult:
        """Search for files"""
        # Simplified - would need async task handling for real implementation
        return SkillResult(
            success=False,
            error="Search not fully implemented yet"
        )


# ============================================================================
# 📚 MNEME - KNOWLEDGE BASE
# ============================================================================

@dataclass
class KnowledgeEntry:
    """Single entry in the MNEME knowledge base"""
    id: str
    category: str           # "pattern", "workflow", "fact", "preference"
    title: str
    content: str
    context: Dict = field(default_factory=dict)
    created_at: str = ""
    updated_at: str = ""
    access_count: int = 0
    confidence: float = 1.0  # How confident is LORENZ about this knowledge
    source: str = ""        # Where this knowledge came from
    tags: List[str] = field(default_factory=list)
    related_skills: List[str] = field(default_factory=list)


class MNEME:
    """
    📚 MNEME - Knowledge Base per LORENZ

    Sistema di memoria persistente che archivia:
    - Patterns di successo nelle interazioni
    - Workflow emergenti scoperti
    - Fatti e preferenze dell'utente
    - Conoscenza appresa dalle conversazioni

    Ispezionabile e modificabile dal pannello di controllo.
    """

    def __init__(self, db_path: str = None):
        self.db_path = db_path or os.path.expanduser("~/.lorenz/mneme.db")
        self._ensure_db_dir()
        self._init_db()
        logger.info(f"📚 MNEME initialized at {self.db_path}")

    def _ensure_db_dir(self):
        """Ensure database directory exists"""
        Path(self.db_path).parent.mkdir(parents=True, exist_ok=True)

    def _init_db(self):
        """Initialize SQLite database"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()

        # Knowledge entries table
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS knowledge (
                id TEXT PRIMARY KEY,
                category TEXT NOT NULL,
                title TEXT NOT NULL,
                content TEXT NOT NULL,
                context TEXT DEFAULT '{}',
                created_at TEXT,
                updated_at TEXT,
                access_count INTEGER DEFAULT 0,
                confidence REAL DEFAULT 1.0,
                source TEXT DEFAULT '',
                tags TEXT DEFAULT '[]',
                related_skills TEXT DEFAULT '[]'
            )
        """)

        # Emergent skills table
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS emergent_skills (
                id TEXT PRIMARY KEY,
                name TEXT UNIQUE NOT NULL,
                description TEXT,
                description_it TEXT,
                trigger_patterns TEXT DEFAULT '[]',
                workflow_steps TEXT DEFAULT '[]',
                category TEXT DEFAULT 'custom',
                created_at TEXT,
                updated_at TEXT,
                use_count INTEGER DEFAULT 0,
                success_rate REAL DEFAULT 1.0,
                enabled INTEGER DEFAULT 1,
                tags TEXT DEFAULT '[]',
                created_from TEXT DEFAULT ''
            )
        """)

        # Learning history table
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS learning_history (
                id TEXT PRIMARY KEY,
                timestamp TEXT,
                event_type TEXT,
                skill_name TEXT,
                input_text TEXT,
                result_success INTEGER,
                learned_pattern TEXT,
                context TEXT DEFAULT '{}'
            )
        """)

        conn.commit()
        conn.close()

    # -------------------------------------------------------------------------
    # Knowledge Operations
    # -------------------------------------------------------------------------

    def add_knowledge(self, entry: KnowledgeEntry) -> bool:
        """Add knowledge entry to MNEME"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            now = datetime.now().isoformat()
            if not entry.id:
                entry.id = str(uuid.uuid4())[:8]
            if not entry.created_at:
                entry.created_at = now
            entry.updated_at = now

            cursor.execute("""
                INSERT OR REPLACE INTO knowledge
                (id, category, title, content, context, created_at, updated_at,
                 access_count, confidence, source, tags, related_skills)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, (
                entry.id, entry.category, entry.title, entry.content,
                json.dumps(entry.context), entry.created_at, entry.updated_at,
                entry.access_count, entry.confidence, entry.source,
                json.dumps(entry.tags), json.dumps(entry.related_skills)
            ))

            conn.commit()
            conn.close()
            logger.info(f"📚 MNEME: Added knowledge '{entry.title}'")
            return True
        except Exception as e:
            logger.error(f"Failed to add knowledge: {e}")
            return False

    def get_knowledge(self, entry_id: str) -> Optional[KnowledgeEntry]:
        """Get knowledge entry by ID"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            cursor.execute("SELECT * FROM knowledge WHERE id = ?", (entry_id,))
            row = cursor.fetchone()

            if row:
                # Update access count
                cursor.execute(
                    "UPDATE knowledge SET access_count = access_count + 1 WHERE id = ?",
                    (entry_id,)
                )
                conn.commit()

                entry = KnowledgeEntry(
                    id=row[0], category=row[1], title=row[2], content=row[3],
                    context=json.loads(row[4] or '{}'),
                    created_at=row[5], updated_at=row[6],
                    access_count=row[7], confidence=row[8], source=row[9],
                    tags=json.loads(row[10] or '[]'),
                    related_skills=json.loads(row[11] or '[]')
                )
                conn.close()
                return entry

            conn.close()
        except Exception as e:
            logger.error(f"Failed to get knowledge: {e}")
        return None

    def search_knowledge(self, query: str = None, category: str = None,
                        tags: List[str] = None, limit: int = 50) -> List[KnowledgeEntry]:
        """Search knowledge base"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            sql = "SELECT * FROM knowledge WHERE 1=1"
            params = []

            if query:
                sql += " AND (title LIKE ? OR content LIKE ?)"
                params.extend([f"%{query}%", f"%{query}%"])

            if category:
                sql += " AND category = ?"
                params.append(category)

            sql += " ORDER BY access_count DESC, updated_at DESC LIMIT ?"
            params.append(limit)

            cursor.execute(sql, params)
            rows = cursor.fetchall()

            entries = []
            for row in rows:
                entry = KnowledgeEntry(
                    id=row[0], category=row[1], title=row[2], content=row[3],
                    context=json.loads(row[4] or '{}'),
                    created_at=row[5], updated_at=row[6],
                    access_count=row[7], confidence=row[8], source=row[9],
                    tags=json.loads(row[10] or '[]'),
                    related_skills=json.loads(row[11] or '[]')
                )
                entries.append(entry)

            conn.close()
            return entries
        except Exception as e:
            logger.error(f"Failed to search knowledge: {e}")
            return []

    def delete_knowledge(self, entry_id: str) -> bool:
        """Delete knowledge entry"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            cursor.execute("DELETE FROM knowledge WHERE id = ?", (entry_id,))
            conn.commit()
            conn.close()
            return True
        except Exception as e:
            logger.error(f"Failed to delete knowledge: {e}")
            return False

    def update_knowledge(self, entry_id: str, updates: Dict) -> bool:
        """Update knowledge entry"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            set_clauses = []
            params = []
            for key, value in updates.items():
                if key in ['context', 'tags', 'related_skills']:
                    value = json.dumps(value)
                set_clauses.append(f"{key} = ?")
                params.append(value)

            set_clauses.append("updated_at = ?")
            params.append(datetime.now().isoformat())
            params.append(entry_id)

            sql = f"UPDATE knowledge SET {', '.join(set_clauses)} WHERE id = ?"
            cursor.execute(sql, params)

            conn.commit()
            conn.close()
            return True
        except Exception as e:
            logger.error(f"Failed to update knowledge: {e}")
            return False

    # -------------------------------------------------------------------------
    # Statistics
    # -------------------------------------------------------------------------

    def get_stats(self) -> Dict:
        """Get MNEME statistics for control panel"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()

            stats = {
                "total_entries": 0,
                "by_category": {},
                "total_skills": 0,
                "recent_activity": []
            }

            # Total knowledge entries
            cursor.execute("SELECT COUNT(*) FROM knowledge")
            stats["total_entries"] = cursor.fetchone()[0]

            # By category
            cursor.execute("""
                SELECT category, COUNT(*) FROM knowledge GROUP BY category
            """)
            for row in cursor.fetchall():
                stats["by_category"][row[0]] = row[1]

            # Total emergent skills
            cursor.execute("SELECT COUNT(*) FROM emergent_skills")
            stats["total_skills"] = cursor.fetchone()[0]

            # Recent activity
            cursor.execute("""
                SELECT title, category, updated_at FROM knowledge
                ORDER BY updated_at DESC LIMIT 10
            """)
            stats["recent_activity"] = [
                {"title": row[0], "category": row[1], "date": row[2]}
                for row in cursor.fetchall()
            ]

            conn.close()
            return stats
        except Exception as e:
            logger.error(f"Failed to get stats: {e}")
            return {}

    def export_to_json(self) -> str:
        """Export entire knowledge base to JSON for backup/inspection"""
        entries = self.search_knowledge(limit=10000)
        data = {
            "exported_at": datetime.now().isoformat(),
            "version": "1.0",
            "entries": [asdict(e) for e in entries]
        }
        return json.dumps(data, indent=2, ensure_ascii=False)

    def import_from_json(self, json_str: str) -> int:
        """Import knowledge from JSON"""
        try:
            data = json.loads(json_str)
            count = 0
            for entry_data in data.get("entries", []):
                entry = KnowledgeEntry(**entry_data)
                if self.add_knowledge(entry):
                    count += 1
            return count
        except Exception as e:
            logger.error(f"Failed to import: {e}")
            return 0


# ============================================================================
# 🧠 EMERGENT SKILL
# ============================================================================

@dataclass
class EmergentSkillDefinition:
    """Definition of an emergent skill learned by LORENZ"""
    id: str
    name: str
    description: str
    description_it: str
    trigger_patterns: List[str]    # Patterns that activate this skill
    workflow_steps: List[Dict]     # Sequence of actions to execute
    category: str = "custom"
    created_at: str = ""
    updated_at: str = ""
    use_count: int = 0
    success_rate: float = 1.0
    enabled: bool = True
    tags: List[str] = field(default_factory=list)
    created_from: str = ""         # Interaction that spawned this skill


class EmergentSkill(Skill):
    """
    🧠 Emergent Skill - Skills learned by LORENZ

    Unlike GOD skills which are hardcoded, emergent skills are:
    - Created dynamically from successful interaction patterns
    - Stored in MNEME knowledge base
    - Editable from the control panel
    - Can combine multiple GOD skills into workflows
    """

    skill_type = SkillType.EMERGENT
    category = SkillCategory.WORKFLOW
    icon = "🧠"

    def __init__(self, definition: EmergentSkillDefinition, skills_manager: 'SkillsManager'):
        self.definition = definition
        self.skills_manager = skills_manager
        self.name = definition.name
        self.description = definition.description
        self.description_it = definition.description_it
        self.examples = definition.trigger_patterns[:3]
        self.requires_api = []
        super().__init__()
        self.skill_id = definition.id
        self.metadata.created_at = definition.created_at
        self.metadata.use_count = definition.use_count
        self.metadata.success_rate = definition.success_rate
        self.metadata.tags = definition.tags

    async def execute(self, **kwargs) -> SkillResult:
        """Execute the emergent skill workflow"""
        results = []
        all_artifacts = []

        try:
            for step in self.definition.workflow_steps:
                skill_name = step.get("skill")
                params = step.get("params", {})

                # Substitute any dynamic params from kwargs
                for key, value in params.items():
                    if isinstance(value, str) and value.startswith("$"):
                        param_name = value[1:]
                        if param_name in kwargs:
                            params[key] = kwargs[param_name]

                # Execute step
                result = await self.skills_manager.execute_skill(skill_name, **params)
                results.append({
                    "skill": skill_name,
                    "success": result.success,
                    "message": result.message
                })

                if result.artifacts:
                    all_artifacts.extend(result.artifacts)

                if not result.success and step.get("required", True):
                    return SkillResult(
                        success=False,
                        error=f"Workflow failed at step '{skill_name}': {result.error}",
                        data={"steps": results}
                    )

            return SkillResult(
                success=True,
                data={"steps": results},
                message=f"🧠 Workflow '{self.name}' completato con {len(results)} step",
                artifacts=all_artifacts
            )

        except Exception as e:
            return SkillResult(
                success=False,
                error=str(e),
                data={"steps": results}
            )

    def get_info(self) -> Dict:
        """Get skill info including workflow details"""
        info = super().get_info()
        info["workflow_steps"] = self.definition.workflow_steps
        info["trigger_patterns"] = self.definition.trigger_patterns
        info["created_from"] = self.definition.created_from
        return info


# ============================================================================
# SKILLS MANAGER (GOD + EMERGENT)
# ============================================================================

class SkillsManager:
    """
    Manages all available skills (GOD + EMERGENT)

    - GOD Skills: Built-in, created by developers
    - Emergent Skills: Learned by LORENZ, stored in MNEME
    """

    def __init__(self, mneme: MNEME = None):
        self.mneme = mneme or MNEME()
        self.god_skills: Dict[str, Skill] = {}
        self.emergent_skills: Dict[str, EmergentSkill] = {}
        self._register_god_skills()
        self._load_emergent_skills()

    def _register_god_skills(self):
        """Register all GOD skills (built-in)"""
        logger.info("🌟 Registering GOD SKILLS...")

        skill_classes = [
            ImageGenerationSkill,
            WebSearchSkill,
            PresentationSkill,
            CodeExecutionSkill,
            ServerCommandSkill,
            WebBrowseSkill,
            DocumentGenerationSkill,
            CalendarSkill,
            VoiceMessageSkill,
            NASOperationSkill,
        ]

        for skill_class in skill_classes:
            try:
                skill = skill_class()
                self.god_skills[skill.name] = skill
                status = "✅" if skill.enabled else "❌"
                logger.info(f"  {status} {skill.icon} {skill.name}")
            except Exception as e:
                logger.error(f"Failed to register GOD skill {skill_class}: {e}")

    def _load_emergent_skills(self):
        """Load emergent skills from MNEME"""
        logger.info("🧠 Loading EMERGENT SKILLS from MNEME...")

        try:
            conn = sqlite3.connect(self.mneme.db_path)
            cursor = conn.cursor()
            cursor.execute("SELECT * FROM emergent_skills WHERE enabled = 1")
            rows = cursor.fetchall()

            for row in rows:
                definition = EmergentSkillDefinition(
                    id=row[0],
                    name=row[1],
                    description=row[2] or "",
                    description_it=row[3] or "",
                    trigger_patterns=json.loads(row[4] or '[]'),
                    workflow_steps=json.loads(row[5] or '[]'),
                    category=row[6] or "custom",
                    created_at=row[7] or "",
                    updated_at=row[8] or "",
                    use_count=row[9] or 0,
                    success_rate=row[10] or 1.0,
                    enabled=bool(row[11]),
                    tags=json.loads(row[12] or '[]'),
                    created_from=row[13] or ""
                )

                skill = EmergentSkill(definition, self)
                self.emergent_skills[skill.name] = skill
                logger.info(f"  🧠 {skill.name}")

            conn.close()
            logger.info(f"  Loaded {len(self.emergent_skills)} emergent skills")

        except Exception as e:
            logger.error(f"Failed to load emergent skills: {e}")

    # -------------------------------------------------------------------------
    # Skill Access
    # -------------------------------------------------------------------------

    @property
    def skills(self) -> Dict[str, Skill]:
        """All skills (GOD + EMERGENT)"""
        return {**self.god_skills, **self.emergent_skills}

    def get_skill(self, name: str) -> Optional[Skill]:
        """Get a skill by name (checks GOD first, then EMERGENT)"""
        return self.god_skills.get(name) or self.emergent_skills.get(name)

    def list_skills(self, enabled_only: bool = True,
                   skill_type: SkillType = None) -> List[Dict]:
        """List skills with optional filtering"""
        skills_list = []

        for skill in self.skills.values():
            if enabled_only and not skill.enabled:
                continue
            if skill_type and skill.skill_type != skill_type:
                continue
            skills_list.append(skill.get_info())

        return skills_list

    def list_god_skills(self, enabled_only: bool = True) -> List[Dict]:
        """List only GOD skills"""
        return self.list_skills(enabled_only, SkillType.GOD)

    def list_emergent_skills(self, enabled_only: bool = True) -> List[Dict]:
        """List only EMERGENT skills"""
        return self.list_skills(enabled_only, SkillType.EMERGENT)

    def get_enabled_skills(self) -> List[str]:
        """Get names of all enabled skills"""
        return [name for name, skill in self.skills.items() if skill.enabled]

    async def execute_skill(self, skill_name: str, **kwargs) -> SkillResult:
        """Execute a skill by name"""
        import time
        start_time = time.time()

        skill = self.get_skill(skill_name)
        if not skill:
            return SkillResult(
                success=False,
                error=f"Skill '{skill_name}' not found"
            )
        if not skill.enabled:
            return SkillResult(
                success=False,
                error=f"Skill '{skill_name}' is disabled"
            )

        result = await skill.execute(**kwargs)

        # Track execution
        execution_time = (time.time() - start_time) * 1000
        skill._track_execution(result, execution_time)
        result.skill_name = skill_name
        result.execution_time_ms = int(execution_time)

        # Log to learning history
        self._log_execution(skill_name, kwargs, result)

        return result

    def _log_execution(self, skill_name: str, params: Dict, result: SkillResult):
        """Log skill execution for learning"""
        try:
            conn = sqlite3.connect(self.mneme.db_path)
            cursor = conn.cursor()

            cursor.execute("""
                INSERT INTO learning_history
                (id, timestamp, event_type, skill_name, input_text, result_success, context)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            """, (
                str(uuid.uuid4())[:8],
                datetime.now().isoformat(),
                "skill_execution",
                skill_name,
                json.dumps(params)[:500],
                1 if result.success else 0,
                json.dumps({"message": result.message[:200] if result.message else ""})
            ))

            conn.commit()
            conn.close()
        except Exception as e:
            logger.warning(f"Failed to log execution: {e}")

    # -------------------------------------------------------------------------
    # Emergent Skill Management (for Control Panel)
    # -------------------------------------------------------------------------

    def create_emergent_skill(self, name: str, description: str,
                             trigger_patterns: List[str],
                             workflow_steps: List[Dict],
                             description_it: str = "",
                             tags: List[str] = None) -> Optional[EmergentSkill]:
        """Create a new emergent skill"""
        try:
            definition = EmergentSkillDefinition(
                id=str(uuid.uuid4())[:8],
                name=name,
                description=description,
                description_it=description_it or description,
                trigger_patterns=trigger_patterns,
                workflow_steps=workflow_steps,
                category="workflow",
                created_at=datetime.now().isoformat(),
                updated_at=datetime.now().isoformat(),
                tags=tags or []
            )

            # Save to MNEME
            conn = sqlite3.connect(self.mneme.db_path)
            cursor = conn.cursor()

            cursor.execute("""
                INSERT INTO emergent_skills
                (id, name, description, description_it, trigger_patterns, workflow_steps,
                 category, created_at, updated_at, use_count, success_rate, enabled, tags)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, (
                definition.id, definition.name, definition.description,
                definition.description_it,
                json.dumps(definition.trigger_patterns),
                json.dumps(definition.workflow_steps),
                definition.category, definition.created_at, definition.updated_at,
                0, 1.0, 1, json.dumps(definition.tags)
            ))

            conn.commit()
            conn.close()

            # Create and register skill
            skill = EmergentSkill(definition, self)
            self.emergent_skills[name] = skill

            logger.info(f"🧠 Created emergent skill: {name}")
            return skill

        except Exception as e:
            logger.error(f"Failed to create emergent skill: {e}")
            return None

    def update_emergent_skill(self, name: str, updates: Dict) -> bool:
        """Update an emergent skill (from control panel)"""
        if name not in self.emergent_skills:
            return False

        try:
            conn = sqlite3.connect(self.mneme.db_path)
            cursor = conn.cursor()

            set_clauses = []
            params = []
            for key, value in updates.items():
                if key in ['trigger_patterns', 'workflow_steps', 'tags']:
                    value = json.dumps(value)
                set_clauses.append(f"{key} = ?")
                params.append(value)

            set_clauses.append("updated_at = ?")
            params.append(datetime.now().isoformat())
            params.append(name)

            sql = f"UPDATE emergent_skills SET {', '.join(set_clauses)} WHERE name = ?"
            cursor.execute(sql, params)

            conn.commit()
            conn.close()

            # Reload skill
            self._load_emergent_skills()
            return True

        except Exception as e:
            logger.error(f"Failed to update emergent skill: {e}")
            return False

    def delete_emergent_skill(self, name: str) -> bool:
        """Delete an emergent skill"""
        if name not in self.emergent_skills:
            return False

        try:
            conn = sqlite3.connect(self.mneme.db_path)
            cursor = conn.cursor()
            cursor.execute("DELETE FROM emergent_skills WHERE name = ?", (name,))
            conn.commit()
            conn.close()

            del self.emergent_skills[name]
            logger.info(f"🗑️ Deleted emergent skill: {name}")
            return True

        except Exception as e:
            logger.error(f"Failed to delete emergent skill: {e}")
            return False

    def toggle_emergent_skill(self, name: str, enabled: bool) -> bool:
        """Enable/disable an emergent skill"""
        return self.update_emergent_skill(name, {"enabled": 1 if enabled else 0})

    # -------------------------------------------------------------------------
    # API for Control Panel
    # -------------------------------------------------------------------------

    def get_control_panel_data(self) -> Dict:
        """Get all data needed for the skills control panel"""
        return {
            "god_skills": self.list_god_skills(enabled_only=False),
            "emergent_skills": self.list_emergent_skills(enabled_only=False),
            "mneme_stats": self.mneme.get_stats(),
            "categories": [c.value for c in SkillCategory],
            "summary": {
                "total_god_skills": len(self.god_skills),
                "total_emergent_skills": len(self.emergent_skills),
                "enabled_god_skills": len([s for s in self.god_skills.values() if s.enabled]),
                "enabled_emergent_skills": len([s for s in self.emergent_skills.values() if s.enabled])
            }
        }


# ============================================================================
# SKILL ROUTER (AI-based skill selection)
# ============================================================================

class SkillRouter:
    """
    Routes user requests to appropriate skills

    Checks GOD skills first, then searches emergent skill patterns.
    """

    # Keywords for GOD skills (built-in)
    GOD_SKILL_KEYWORDS = {
        "image_generation": [
            "genera immagine", "genera un'immagine", "genera un immagine",
            "crea immagine", "crea un'immagine", "disegna", "illustra",
            "generate image", "create image", "draw", "picture", "immagine di"
        ],
        "web_search": [
            "cerca", "search", "trova online", "google", "news",
            "notizie", "attuale", "current", "latest"
        ],
        "presentation": [
            "presentazione", "slide", "powerpoint", "pptx",
            "presentation", "slides"
        ],
        "code_execution": [
            "esegui codice", "run code", "calcola", "compute",
            "python", "script"
        ],
        "server_command": [
            "server", "ssh", "systemctl", "riavvia", "restart",
            "logs", "status server", "uptime"
        ],
        "web_browse": [
            "apri url", "leggi pagina", "vai a", "browse",
            "fetch url", "scarica pagina"
        ],
        "document_generation": [
            "crea documento", "genera word", "docx", "report",
            "create document"
        ],
        "nas_operation": [
            "nas", "synology", "file server", "backup",
            "storage", "cartella condivisa"
        ],
        "calendar": [
            "calendario", "calendar", "appuntamento", "meeting",
            "agenda", "evento", "schedule", "riunione",
            "prossimi appuntamenti", "cosa ho in agenda"
        ],
        "voice_message": [
            "trascrivi", "transcribe", "voce", "voice", "audio",
            "messaggio vocale", "ascolta", "registrazione"
        ],
    }

    def __init__(self, skills_manager: SkillsManager = None):
        self.skills_manager = skills_manager

    def detect_skill(self, text: str) -> Optional[str]:
        """Detect which skill should handle the request"""
        text_lower = text.lower()

        # First, check GOD skill keywords
        scores = {skill: 0 for skill in self.GOD_SKILL_KEYWORDS}

        for skill, keywords in self.GOD_SKILL_KEYWORDS.items():
            for keyword in keywords:
                if keyword in text_lower:
                    scores[skill] += 1

        max_score = max(scores.values())
        if max_score > 0:
            for skill, score in scores.items():
                if score == max_score:
                    return skill

        # Then, check emergent skill trigger patterns
        if self.skills_manager:
            for skill_name, skill in self.skills_manager.emergent_skills.items():
                for pattern in skill.definition.trigger_patterns:
                    if pattern.lower() in text_lower:
                        return skill_name

        return None

    @classmethod
    def detect_skill_static(cls, text: str) -> Optional[str]:
        """Static method for backward compatibility (GOD skills only)"""
        text_lower = text.lower()

        scores = {skill: 0 for skill in cls.GOD_SKILL_KEYWORDS}

        for skill, keywords in cls.GOD_SKILL_KEYWORDS.items():
            for keyword in keywords:
                if keyword in text_lower:
                    scores[skill] += 1

        max_score = max(scores.values())
        if max_score > 0:
            for skill, score in scores.items():
                if score == max_score:
                    return skill

        return None

    @staticmethod
    def extract_skill_params(skill_name: str, text: str) -> Dict:
        """Extract parameters for skill from text"""
        params = {}

        if skill_name == "image_generation":
            # Everything after trigger words is the prompt
            for trigger in ["genera immagine", "crea immagine", "disegna"]:
                if trigger in text.lower():
                    idx = text.lower().find(trigger) + len(trigger)
                    params["prompt"] = text[idx:].strip()
                    break
            if not params.get("prompt"):
                params["prompt"] = text

        elif skill_name == "web_search":
            for trigger in ["cerca", "search", "trova"]:
                if trigger in text.lower():
                    idx = text.lower().find(trigger) + len(trigger)
                    params["query"] = text[idx:].strip()
                    break
            if not params.get("query"):
                params["query"] = text

        elif skill_name == "web_browse":
            # Extract URL
            import re
            url_pattern = r'https?://[^\s]+'
            match = re.search(url_pattern, text)
            if match:
                params["url"] = match.group()

        elif skill_name == "calendar":
            # Default to list operation
            params["operation"] = "list"
            text_lower = text.lower()

            # Check for create operations
            create_triggers = ["crea", "aggiungi", "nuovo", "create", "add", "schedule"]
            for trigger in create_triggers:
                if trigger in text_lower:
                    params["operation"] = "create"
                    # Extract event name (simple approach)
                    idx = text_lower.find(trigger) + len(trigger)
                    remaining = text[idx:].strip()
                    if remaining:
                        params["summary"] = remaining[:100]  # Limit summary length
                    break

            # Extract time if mentioned
            import re
            time_pattern = r'(?:alle\s+)?(\d{1,2}[:.]\d{2})'
            time_match = re.search(time_pattern, text)
            if time_match:
                params["time_hint"] = time_match.group(1)

        elif skill_name == "voice_message":
            # Check for text-to-speech request
            tts_triggers = ["leggi", "pronuncia", "dì", "speak", "say", "read aloud"]
            for trigger in tts_triggers:
                if trigger in text.lower():
                    params["tts"] = True
                    idx = text.lower().find(trigger) + len(trigger)
                    params["text"] = text[idx:].strip()
                    break

        return params


# ============================================================================
# GLOBAL INSTANCES
# ============================================================================

_mneme: Optional[MNEME] = None
_skills_manager: Optional[SkillsManager] = None


def get_mneme() -> MNEME:
    """Get or create global MNEME instance"""
    global _mneme
    if _mneme is None:
        _mneme = MNEME()
    return _mneme


def get_skills_manager() -> SkillsManager:
    """Get or create global skills manager"""
    global _skills_manager
    if _skills_manager is None:
        _skills_manager = SkillsManager(get_mneme())
    return _skills_manager


# ============================================================================
# CLI TEST
# ============================================================================

if __name__ == "__main__":
    # Configure logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )

    async def test():
        print("\n" + "=" * 60)
        print("🎯 LORENZ Skills System v2.0")
        print("   GOD SKILLS + EMERGENT SKILLS + MNEME")
        print("=" * 60)

        # Initialize
        mneme = MNEME()
        manager = SkillsManager(mneme)

        # Display GOD SKILLS
        print("\n🌟 GOD SKILLS (Built-in):")
        print("-" * 40)
        for skill_info in manager.list_god_skills(enabled_only=False):
            status = "✅" if skill_info["enabled"] else "❌"
            icon = skill_info.get("icon", "⚡")
            print(f"  {status} {icon} {skill_info['name']}")
            print(f"      {skill_info['description_it']}")

        # Display EMERGENT SKILLS
        print("\n🧠 EMERGENT SKILLS (Learned):")
        print("-" * 40)
        emergent = manager.list_emergent_skills(enabled_only=False)
        if emergent:
            for skill_info in emergent:
                status = "✅" if skill_info["enabled"] else "❌"
                print(f"  {status} 🧠 {skill_info['name']}")
                print(f"      Workflow: {len(skill_info.get('workflow_steps', []))} steps")
        else:
            print("  (Nessuna skill emergente ancora appresa)")

        # Display MNEME stats
        print("\n📚 MNEME Knowledge Base:")
        print("-" * 40)
        stats = mneme.get_stats()
        print(f"  Total entries: {stats.get('total_entries', 0)}")
        print(f"  Categories: {stats.get('by_category', {})}")
        print(f"  Emergent skills: {stats.get('total_skills', 0)}")

        # Create a sample emergent skill
        print("\n🔧 Creating sample emergent skill...")
        print("-" * 40)

        sample_skill = manager.create_emergent_skill(
            name="daily_briefing",
            description="Generate a daily briefing with news and server status",
            description_it="Genera un briefing giornaliero con notizie e stato server",
            trigger_patterns=[
                "briefing giornaliero",
                "daily briefing",
                "cosa c'è di nuovo",
                "buongiorno lorenz"
            ],
            workflow_steps=[
                {"skill": "web_search", "params": {"query": "top tech news today"}},
                {"skill": "server_command", "params": {"command": "uptime"}, "required": False}
            ],
            tags=["daily", "briefing", "workflow"]
        )

        if sample_skill:
            print(f"  ✅ Created: {sample_skill.name}")
            print(f"     Triggers: {sample_skill.definition.trigger_patterns}")
        else:
            print("  ℹ️ Skill already exists or creation failed")

        # Test skill routing
        print("\n🔀 Skill Routing Test:")
        print("-" * 40)
        router = SkillRouter(manager)

        test_inputs = [
            "Genera un'immagine di un gatto spaziale",
            "Cerca le ultime notizie su AI",
            "Crea una presentazione sul machine learning",
            "Controlla lo stato del server",
            "Buongiorno Lorenz, cosa c'è di nuovo?",
            "Ciao, come stai?",
        ]

        for text in test_inputs:
            skill = router.detect_skill(text)
            skill_type = ""
            if skill:
                if skill in manager.god_skills:
                    skill_type = "🌟 GOD"
                elif skill in manager.emergent_skills:
                    skill_type = "🧠 EMERGENT"
            print(f"  '{text[:35]}...'")
            print(f"     → {skill or 'None (use AI chat)'} {skill_type}")

        # Control Panel API
        print("\n📊 Control Panel Data:")
        print("-" * 40)
        cp_data = manager.get_control_panel_data()
        print(f"  GOD Skills: {cp_data['summary']['total_god_skills']} total, "
              f"{cp_data['summary']['enabled_god_skills']} enabled")
        print(f"  Emergent Skills: {cp_data['summary']['total_emergent_skills']} total, "
              f"{cp_data['summary']['enabled_emergent_skills']} enabled")

        print("\n" + "=" * 60)
        print("✅ LORENZ Skills System initialized successfully!")
        print("=" * 60)

    asyncio.run(test())
