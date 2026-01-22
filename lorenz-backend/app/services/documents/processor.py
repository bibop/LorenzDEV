"""
LORENZ SaaS - Document Processor
=================================

Advanced document processing for RAG indexing with:
- Multi-format text extraction
- Smart chunking with overlap
- Metadata extraction
- OCR support for images/scanned PDFs
- Async processing with progress tracking
"""

import os
import io
import re
import hashlib
import logging
import asyncio
from typing import List, Dict, Optional, Any, Tuple, BinaryIO
from dataclasses import dataclass, field
from enum import Enum
from datetime import datetime
from uuid import UUID, uuid4
from concurrent.futures import ThreadPoolExecutor

logger = logging.getLogger(__name__)

# Thread pool for CPU-bound operations
_executor = ThreadPoolExecutor(max_workers=4)


# ============================================================================
# ENUMS AND DATA CLASSES
# ============================================================================

class ProcessingStatus(str, Enum):
    """Document processing status"""
    PENDING = "pending"
    EXTRACTING = "extracting"
    CHUNKING = "chunking"
    INDEXING = "indexing"
    COMPLETED = "completed"
    FAILED = "failed"


class ChunkingStrategy(str, Enum):
    """Text chunking strategies"""
    FIXED = "fixed"              # Fixed size chunks
    SENTENCE = "sentence"        # Sentence-based
    PARAGRAPH = "paragraph"      # Paragraph-based
    SEMANTIC = "semantic"        # Semantic similarity (advanced)
    RECURSIVE = "recursive"      # Recursive splitting


@dataclass
class DocumentChunk:
    """A chunk of document text"""
    id: str
    document_id: str
    content: str
    chunk_index: int
    total_chunks: int
    start_char: int
    end_char: int
    metadata: Dict[str, Any] = field(default_factory=dict)
    embedding: Optional[List[float]] = None

    def to_dict(self) -> Dict:
        return {
            "id": self.id,
            "document_id": self.document_id,
            "content": self.content,
            "chunk_index": self.chunk_index,
            "total_chunks": self.total_chunks,
            "start_char": self.start_char,
            "end_char": self.end_char,
            "metadata": self.metadata,
        }


@dataclass
class ProcessingResult:
    """Result of document processing"""
    document_id: str
    filename: str
    status: ProcessingStatus
    content_hash: str
    text_length: int
    chunks: List[DocumentChunk] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    error: Optional[str] = None
    processing_time_ms: int = 0

    def to_dict(self) -> Dict:
        return {
            "document_id": self.document_id,
            "filename": self.filename,
            "status": self.status.value,
            "content_hash": self.content_hash,
            "text_length": self.text_length,
            "chunk_count": len(self.chunks),
            "metadata": self.metadata,
            "error": self.error,
            "processing_time_ms": self.processing_time_ms,
        }


# ============================================================================
# SUPPORTED MIME TYPES
# ============================================================================

MIME_TYPE_HANDLERS = {
    # Text formats
    "text/plain": "extract_text_plain",
    "text/markdown": "extract_text_plain",
    "text/html": "extract_html",
    "text/csv": "extract_csv",

    # PDF
    "application/pdf": "extract_pdf",

    # Microsoft Office
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "extract_docx",
    "application/msword": "extract_doc",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "extract_xlsx",
    "application/vnd.ms-excel": "extract_xls",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "extract_pptx",

    # LibreOffice/OpenOffice
    "application/vnd.oasis.opendocument.text": "extract_odt",

    # Images (OCR)
    "image/png": "extract_image_ocr",
    "image/jpeg": "extract_image_ocr",
    "image/tiff": "extract_image_ocr",
    "image/webp": "extract_image_ocr",

    # Rich Text
    "application/rtf": "extract_rtf",
}

# File extension to MIME type mapping
EXTENSION_MIME_MAP = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".html": "text/html",
    ".htm": "text/html",
    ".csv": "text/csv",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls": "application/vnd.ms-excel",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".odt": "application/vnd.oasis.opendocument.text",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
    ".webp": "image/webp",
    ".rtf": "application/rtf",
}


# ============================================================================
# DOCUMENT PROCESSOR CLASS
# ============================================================================

class DocumentProcessor:
    """
    Multi-format document processor for RAG indexing

    Features:
    - Automatic format detection
    - Text extraction from multiple formats
    - Smart chunking with configurable strategies
    - Metadata extraction
    - OCR for images and scanned PDFs
    """

    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        chunking_strategy: ChunkingStrategy = ChunkingStrategy.RECURSIVE,
        enable_ocr: bool = True,
        ocr_language: str = "eng+ita"
    ):
        """
        Initialize document processor

        Args:
            chunk_size: Target chunk size in characters
            chunk_overlap: Overlap between chunks
            chunking_strategy: Strategy for splitting text
            enable_ocr: Enable OCR for images
            ocr_language: Tesseract language codes
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.chunking_strategy = chunking_strategy
        self.enable_ocr = enable_ocr
        self.ocr_language = ocr_language

        # Processing callbacks
        self._progress_callback = None

    def set_progress_callback(self, callback):
        """Set callback for progress updates: callback(status, progress_pct)"""
        self._progress_callback = callback

    def _update_progress(self, status: ProcessingStatus, progress: float = 0):
        """Update processing progress"""
        if self._progress_callback:
            try:
                self._progress_callback(status, progress)
            except Exception:
                pass

    def detect_mime_type(self, filename: str, content: Optional[bytes] = None) -> str:
        """Detect MIME type from filename or content"""
        # Try by extension first
        ext = os.path.splitext(filename.lower())[1]
        if ext in EXTENSION_MIME_MAP:
            return EXTENSION_MIME_MAP[ext]

        # Try magic number detection
        if content:
            if content[:4] == b'%PDF':
                return "application/pdf"
            if content[:4] == b'PK\x03\x04':
                # Could be DOCX, XLSX, PPTX, etc.
                if b'word/' in content[:2000]:
                    return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                if b'xl/' in content[:2000]:
                    return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            if content[:8] == b'\x89PNG\r\n\x1a\n':
                return "image/png"
            if content[:2] == b'\xff\xd8':
                return "image/jpeg"

        return "application/octet-stream"

    async def process_document(
        self,
        content: bytes,
        filename: str,
        content_type: Optional[str] = None,
        metadata: Optional[Dict] = None
    ) -> ProcessingResult:
        """
        Process a document: extract text and create chunks

        Args:
            content: Document binary content
            filename: Original filename
            content_type: MIME type (auto-detected if not provided)
            metadata: Additional metadata to attach

        Returns:
            ProcessingResult with chunks
        """
        start_time = datetime.now()
        document_id = str(uuid4())

        self._update_progress(ProcessingStatus.EXTRACTING, 0)

        try:
            # Detect content type
            if not content_type or content_type == "application/octet-stream":
                content_type = self.detect_mime_type(filename, content)

            logger.info(f"Processing {filename} ({content_type})")

            # Get handler
            handler_name = MIME_TYPE_HANDLERS.get(content_type)
            if not handler_name:
                raise ValueError(f"Unsupported content type: {content_type}")

            handler = getattr(self, handler_name, None)
            if not handler:
                raise ValueError(f"Handler not implemented: {handler_name}")

            # Extract text (run in executor for CPU-bound operations)
            loop = asyncio.get_event_loop()
            text, doc_metadata = await loop.run_in_executor(
                _executor,
                lambda: handler(content, filename)
            )

            self._update_progress(ProcessingStatus.EXTRACTING, 50)

            if not text or len(text.strip()) < 10:
                raise ValueError("No text extracted from document")

            # Clean text
            text = self._clean_text(text)

            # Compute content hash
            content_hash = hashlib.sha256(text.encode()).hexdigest()

            self._update_progress(ProcessingStatus.CHUNKING, 60)

            # Create chunks
            chunks = await self._create_chunks(
                text=text,
                document_id=document_id,
                filename=filename
            )

            self._update_progress(ProcessingStatus.COMPLETED, 100)

            # Calculate processing time
            processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

            # Merge metadata
            final_metadata = {
                "filename": filename,
                "content_type": content_type,
                "original_size": len(content),
                "extracted_at": datetime.now().isoformat(),
                **(doc_metadata or {}),
                **(metadata or {})
            }

            return ProcessingResult(
                document_id=document_id,
                filename=filename,
                status=ProcessingStatus.COMPLETED,
                content_hash=content_hash,
                text_length=len(text),
                chunks=chunks,
                metadata=final_metadata,
                processing_time_ms=processing_time
            )

        except Exception as e:
            logger.error(f"Document processing failed: {e}", exc_info=True)
            return ProcessingResult(
                document_id=document_id,
                filename=filename,
                status=ProcessingStatus.FAILED,
                content_hash="",
                text_length=0,
                error=str(e),
                processing_time_ms=int((datetime.now() - start_time).total_seconds() * 1000)
            )

    # ========================================================================
    # TEXT EXTRACTION METHODS
    # ========================================================================

    def extract_text_plain(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from plain text files"""
        # Try different encodings
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']

        for encoding in encodings:
            try:
                text = content.decode(encoding)
                return text, {"encoding": encoding}
            except UnicodeDecodeError:
                continue

        # Fallback with errors='replace'
        text = content.decode('utf-8', errors='replace')
        return text, {"encoding": "utf-8-replaced"}

    def extract_html(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from HTML"""
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(content, 'html.parser')

            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()

            # Get text
            text = soup.get_text(separator='\n')

            # Extract metadata
            title = soup.title.string if soup.title else None
            meta_desc = soup.find('meta', attrs={'name': 'description'})

            metadata = {
                "title": title,
                "description": meta_desc.get('content') if meta_desc else None
            }

            return text, metadata

        except ImportError:
            # Fallback: simple tag stripping
            text = re.sub(r'<[^>]+>', '', content.decode('utf-8', errors='replace'))
            return text, {}

    def extract_csv(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from CSV"""
        import csv

        text = content.decode('utf-8', errors='replace')
        reader = csv.reader(io.StringIO(text))

        rows = list(reader)
        if not rows:
            return "", {"rows": 0, "columns": 0}

        # Convert to readable text
        result_lines = []
        headers = rows[0] if rows else []

        for i, row in enumerate(rows):
            if i == 0:
                result_lines.append("Headers: " + ", ".join(row))
            else:
                # Create key-value pairs
                pairs = [f"{headers[j] if j < len(headers) else f'col{j}'}: {cell}"
                        for j, cell in enumerate(row)]
                result_lines.append(f"Row {i}: " + "; ".join(pairs))

        return "\n".join(result_lines), {
            "rows": len(rows),
            "columns": len(headers)
        }

    def extract_pdf(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from PDF (with OCR fallback)"""
        text = ""
        metadata = {}

        try:
            # Try PyPDF2 first (fast, text-based PDFs)
            import PyPDF2

            reader = PyPDF2.PdfReader(io.BytesIO(content))

            pages_text = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    pages_text.append(f"--- Page {i+1} ---\n{page_text}")

            text = "\n\n".join(pages_text)

            # Get metadata
            if reader.metadata:
                metadata = {
                    "title": reader.metadata.get("/Title"),
                    "author": reader.metadata.get("/Author"),
                    "subject": reader.metadata.get("/Subject"),
                    "pages": len(reader.pages)
                }

            # If no text found, try OCR
            if len(text.strip()) < 100 and self.enable_ocr:
                logger.info(f"PDF appears scanned, trying OCR for {filename}")
                text, ocr_meta = self._ocr_pdf(content)
                metadata.update(ocr_meta)
                metadata["ocr_used"] = True

        except ImportError:
            logger.warning("PyPDF2 not installed, trying OCR")
            if self.enable_ocr:
                text, metadata = self._ocr_pdf(content)
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            if self.enable_ocr:
                text, metadata = self._ocr_pdf(content)

        return text, metadata

    def _ocr_pdf(self, content: bytes) -> Tuple[str, Dict]:
        """OCR a PDF document"""
        try:
            import pdf2image
            import pytesseract

            images = pdf2image.convert_from_bytes(content)

            text_parts = []
            for i, image in enumerate(images):
                page_text = pytesseract.image_to_string(
                    image,
                    lang=self.ocr_language
                )
                if page_text.strip():
                    text_parts.append(f"--- Page {i+1} (OCR) ---\n{page_text}")

            return "\n\n".join(text_parts), {"pages": len(images), "ocr": True}

        except ImportError as e:
            logger.warning(f"OCR dependencies not available: {e}")
            return "", {"error": "OCR not available"}

    def extract_docx(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from DOCX"""
        try:
            import docx

            doc = docx.Document(io.BytesIO(content))

            text_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    text_parts.append("\n[Table]\n" + "\n".join(table_text))

            # Get metadata
            metadata = {}
            if doc.core_properties:
                props = doc.core_properties
                metadata = {
                    "title": props.title,
                    "author": props.author,
                    "subject": props.subject,
                    "created": props.created.isoformat() if props.created else None,
                }

            return "\n\n".join(text_parts), metadata

        except ImportError:
            logger.warning("python-docx not installed")
            return "", {"error": "DOCX support not available"}

    def extract_doc(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from DOC (legacy Word)"""
        try:
            import antiword
            text = antiword.parse(io.BytesIO(content))
            return text, {"format": "doc"}
        except ImportError:
            # Try textract as fallback
            try:
                import textract
                text = textract.process(io.BytesIO(content)).decode('utf-8')
                return text, {"format": "doc"}
            except ImportError:
                logger.warning("No DOC parser available (antiword/textract)")
                return "", {"error": "DOC support not available"}

    def extract_xlsx(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from Excel XLSX"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)

            text_parts = []
            sheet_count = 0

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_count += 1

                rows_text = [f"\n=== Sheet: {sheet_name} ==="]

                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        rows_text.append(" | ".join(row_values))

                if len(rows_text) > 1:
                    text_parts.append("\n".join(rows_text))

            return "\n".join(text_parts), {"sheets": sheet_count}

        except ImportError:
            logger.warning("openpyxl not installed")
            return "", {"error": "XLSX support not available"}

    def extract_xls(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from Excel XLS (legacy)"""
        try:
            import xlrd

            wb = xlrd.open_workbook(file_contents=content)

            text_parts = []

            for sheet in wb.sheets():
                rows_text = [f"\n=== Sheet: {sheet.name} ==="]

                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    row_values = [str(cell) for cell in row]
                    if any(v.strip() for v in row_values):
                        rows_text.append(" | ".join(row_values))

                if len(rows_text) > 1:
                    text_parts.append("\n".join(rows_text))

            return "\n".join(text_parts), {"sheets": wb.nsheets}

        except ImportError:
            logger.warning("xlrd not installed")
            return "", {"error": "XLS support not available"}

    def extract_pptx(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from PowerPoint PPTX"""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))

            text_parts = []
            slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_count += 1
                slide_text = [f"\n--- Slide {slide_num} ---"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))

            return "\n".join(text_parts), {"slides": slide_count}

        except ImportError:
            logger.warning("python-pptx not installed")
            return "", {"error": "PPTX support not available"}

    def extract_odt(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from OpenDocument Text"""
        try:
            from odf import text as odf_text
            from odf.opendocument import load

            doc = load(io.BytesIO(content))

            text_parts = []
            for para in doc.getElementsByType(odf_text.P):
                para_text = "".join(
                    node.data for node in para.childNodes
                    if hasattr(node, 'data')
                )
                if para_text.strip():
                    text_parts.append(para_text)

            return "\n\n".join(text_parts), {"format": "odt"}

        except ImportError:
            logger.warning("odfpy not installed")
            return "", {"error": "ODT support not available"}

    def extract_image_ocr(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from images using OCR"""
        if not self.enable_ocr:
            return "", {"error": "OCR disabled"}

        try:
            from PIL import Image
            import pytesseract

            image = Image.open(io.BytesIO(content))

            text = pytesseract.image_to_string(
                image,
                lang=self.ocr_language
            )

            metadata = {
                "width": image.width,
                "height": image.height,
                "format": image.format,
                "ocr": True
            }

            return text, metadata

        except ImportError as e:
            logger.warning(f"OCR dependencies not available: {e}")
            return "", {"error": "OCR not available"}

    def extract_rtf(self, content: bytes, filename: str) -> Tuple[str, Dict]:
        """Extract text from RTF"""
        try:
            from striprtf.striprtf import rtf_to_text

            rtf_content = content.decode('utf-8', errors='replace')
            text = rtf_to_text(rtf_content)

            return text, {"format": "rtf"}

        except ImportError:
            logger.warning("striprtf not installed")
            return "", {"error": "RTF support not available"}

    # ========================================================================
    # TEXT CLEANING AND CHUNKING
    # ========================================================================

    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        # Remove null bytes
        text = text.replace('\x00', '')

        # Normalize whitespace
        text = re.sub(r'[ \t]+', ' ', text)

        # Normalize line breaks
        text = re.sub(r'\n{3,}', '\n\n', text)

        # Remove control characters (except newline, tab)
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)

        return text.strip()

    async def _create_chunks(
        self,
        text: str,
        document_id: str,
        filename: str
    ) -> List[DocumentChunk]:
        """
        Create text chunks using the configured strategy
        """
        if self.chunking_strategy == ChunkingStrategy.RECURSIVE:
            chunks = self._recursive_chunk(text)
        elif self.chunking_strategy == ChunkingStrategy.SENTENCE:
            chunks = self._sentence_chunk(text)
        elif self.chunking_strategy == ChunkingStrategy.PARAGRAPH:
            chunks = self._paragraph_chunk(text)
        else:
            chunks = self._fixed_chunk(text)

        # Convert to DocumentChunk objects
        result = []
        total_chunks = len(chunks)

        for i, (chunk_text, start, end) in enumerate(chunks):
            result.append(DocumentChunk(
                id=f"{document_id}_chunk_{i}",
                document_id=document_id,
                content=chunk_text,
                chunk_index=i,
                total_chunks=total_chunks,
                start_char=start,
                end_char=end,
                metadata={
                    "filename": filename,
                    "chunk_strategy": self.chunking_strategy.value
                }
            ))

        return result

    def _fixed_chunk(self, text: str) -> List[Tuple[str, int, int]]:
        """Simple fixed-size chunking with overlap"""
        chunks = []
        start = 0

        while start < len(text):
            end = min(start + self.chunk_size, len(text))

            # Try to end at word boundary
            if end < len(text):
                last_space = text.rfind(' ', start, end)
                if last_space > start:
                    end = last_space

            chunk = text[start:end].strip()
            if chunk:
                chunks.append((chunk, start, end))

            # Move start with overlap
            start = end - self.chunk_overlap
            if start <= chunks[-1][1] if chunks else 0:
                start = end

        return chunks

    def _sentence_chunk(self, text: str) -> List[Tuple[str, int, int]]:
        """Chunk by sentences, respecting size limits"""
        # Simple sentence splitting
        sentence_endings = re.compile(r'(?<=[.!?])\s+')
        sentences = sentence_endings.split(text)

        chunks = []
        current_chunk = ""
        current_start = 0
        position = 0

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            if len(current_chunk) + len(sentence) + 1 > self.chunk_size:
                if current_chunk:
                    chunks.append((current_chunk, current_start, position))
                current_chunk = sentence
                current_start = position
            else:
                if current_chunk:
                    current_chunk += " " + sentence
                else:
                    current_chunk = sentence
                    current_start = position

            position += len(sentence) + 1

        if current_chunk:
            chunks.append((current_chunk, current_start, len(text)))

        return chunks

    def _paragraph_chunk(self, text: str) -> List[Tuple[str, int, int]]:
        """Chunk by paragraphs, respecting size limits"""
        paragraphs = text.split('\n\n')

        chunks = []
        current_chunk = ""
        current_start = 0
        position = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                position += 2
                continue

            if len(current_chunk) + len(para) + 2 > self.chunk_size:
                if current_chunk:
                    chunks.append((current_chunk, current_start, position))

                # If paragraph itself is too long, split it
                if len(para) > self.chunk_size:
                    sub_chunks = self._fixed_chunk(para)
                    for sub_chunk, sub_start, sub_end in sub_chunks:
                        chunks.append((sub_chunk, position + sub_start, position + sub_end))
                    current_chunk = ""
                else:
                    current_chunk = para
                    current_start = position
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
                    current_start = position

            position += len(para) + 2

        if current_chunk:
            chunks.append((current_chunk, current_start, len(text)))

        return chunks

    def _recursive_chunk(self, text: str) -> List[Tuple[str, int, int]]:
        """
        Recursive chunking: try progressively smaller separators
        """
        separators = [
            "\n\n\n",      # Multiple newlines (sections)
            "\n\n",        # Paragraphs
            "\n",          # Lines
            ". ",          # Sentences
            ", ",          # Clauses
            " ",           # Words
            ""             # Characters
        ]

        return self._recursive_split(text, separators, 0)

    def _recursive_split(
        self,
        text: str,
        separators: List[str],
        offset: int
    ) -> List[Tuple[str, int, int]]:
        """Recursively split text"""
        final_chunks = []

        # Find the best separator
        separator = separators[0] if separators else ""

        if separator:
            splits = text.split(separator)
        else:
            splits = list(text)

        current_chunk = ""
        current_start = offset
        position = offset

        for i, split in enumerate(splits):
            split_with_sep = split + (separator if i < len(splits) - 1 else "")

            if len(current_chunk) + len(split_with_sep) > self.chunk_size:
                # Current chunk is full
                if current_chunk:
                    # Check if we can split further
                    if len(current_chunk) > self.chunk_size and len(separators) > 1:
                        sub_chunks = self._recursive_split(
                            current_chunk,
                            separators[1:],
                            current_start
                        )
                        final_chunks.extend(sub_chunks)
                    else:
                        final_chunks.append((current_chunk, current_start, position))

                # Handle overlap
                if self.chunk_overlap > 0 and current_chunk:
                    overlap_text = current_chunk[-self.chunk_overlap:]
                    current_chunk = overlap_text + split_with_sep
                    current_start = position - len(overlap_text)
                else:
                    current_chunk = split_with_sep
                    current_start = position
            else:
                current_chunk += split_with_sep

            position += len(split_with_sep)

        if current_chunk:
            final_chunks.append((current_chunk, current_start, offset + len(text)))

        return final_chunks


# ============================================================================
# FACTORY FUNCTION
# ============================================================================

def get_document_processor(
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
    chunking_strategy: ChunkingStrategy = ChunkingStrategy.RECURSIVE,
    enable_ocr: bool = True
) -> DocumentProcessor:
    """
    Factory function to create a DocumentProcessor

    Args:
        chunk_size: Target chunk size in characters
        chunk_overlap: Overlap between chunks
        chunking_strategy: How to split text
        enable_ocr: Enable OCR for images/scanned PDFs

    Returns:
        Configured DocumentProcessor instance
    """
    return DocumentProcessor(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        chunking_strategy=chunking_strategy,
        enable_ocr=enable_ocr
    )
