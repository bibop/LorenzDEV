"""
LORENZ SaaS - GOD Skills (Built-in Skills)
==========================================

Built-in capabilities adapted from lorenz_skills.py:
- Image Generation (DALL-E 3)
- Web Search (Perplexity)
- Presentation Creation
- Document Generation
- Code Analysis
- Email Draft
- Calendar Management
"""

import os
import logging
import aiohttp
import tempfile
from typing import Dict, List, Optional, Any
from datetime import datetime

from .base import Skill, SkillResult, SkillType, SkillCategory

logger = logging.getLogger(__name__)


# ============================================================================
# IMAGE GENERATION SKILL
# ============================================================================

class ImageGenerationSkill(Skill):
    """Generate images using DALL-E 3"""

    name = "image_generation"
    description = "Generate images from text descriptions using DALL-E 3"
    description_it = "Genera immagini da descrizioni testuali usando DALL-E 3"
    examples = [
        "Genera un'immagine di un tramonto sul mare",
        "Crea un logo per una startup tech",
        "Disegna un gatto astronauta",
        "Create an image of a futuristic city",
        "Design a book cover for a sci-fi novel"
    ]
    requires_api = ["OPENAI"]
    skill_type = SkillType.GOD
    category = SkillCategory.CREATIVE
    icon = "🖼️"
    estimated_cost_usd = 0.04  # Standard quality

    async def execute(
        self,
        prompt: str,
        size: str = "1024x1024",
        quality: str = "standard",
        style: str = "vivid"
    ) -> SkillResult:
        """
        Generate an image

        Args:
            prompt: Image description
            size: "1024x1024", "1792x1024", "1024x1792"
            quality: "standard" or "hd"
            style: "vivid" or "natural"

        Returns:
            SkillResult with image URL
        """
        start_time = datetime.now()

        if not self.enabled:
            return SkillResult(
                success=False,
                error="OpenAI API key not configured",
                skill_name=self.name
            )

        api_key = self._get_api_key("OPENAI")

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }

        payload = {
            "model": "dall-e-3",
            "prompt": prompt,
            "size": size,
            "quality": quality,
            "style": style,
            "n": 1
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.openai.com/v1/images/generations",
                    headers=headers,
                    json=payload
                ) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        image_url = data["data"][0]["url"]
                        revised_prompt = data["data"][0].get("revised_prompt", prompt)

                        execution_time = (datetime.now() - start_time).total_seconds() * 1000

                        result = SkillResult(
                            success=True,
                            data={
                                "url": image_url,
                                "revised_prompt": revised_prompt,
                                "size": size,
                                "quality": quality
                            },
                            message=f"Image generated!\n\nPrompt: {revised_prompt}",
                            artifacts=[image_url],
                            skill_name=self.name,
                            execution_time_ms=int(execution_time)
                        )

                        self._track_execution(result, execution_time)
                        return result
                    else:
                        error = await resp.text()
                        return SkillResult(
                            success=False,
                            error=f"DALL-E API error: {error}",
                            skill_name=self.name
                        )

        except Exception as e:
            return SkillResult(
                success=False,
                error=str(e),
                skill_name=self.name
            )


# ============================================================================
# WEB SEARCH SKILL
# ============================================================================

class WebSearchSkill(Skill):
    """Search the web using Perplexity"""

    name = "web_search"
    description = "Search for up-to-date information on the web"
    description_it = "Cerca informazioni aggiornate sul web"
    examples = [
        "Cerca le ultime notizie su AI",
        "Trova il prezzo attuale di Bitcoin",
        "Search for latest news about Tesla",
        "Find reviews of iPhone 16"
    ]
    requires_api = ["PERPLEXITY"]
    skill_type = SkillType.GOD
    category = SkillCategory.RESEARCH
    icon = "🔍"
    estimated_cost_usd = 0.005

    async def execute(
        self,
        query: str,
        detailed: bool = False
    ) -> SkillResult:
        """
        Search the web

        Args:
            query: Search query
            detailed: If True, use larger model for detailed results

        Returns:
            SkillResult with search results
        """
        start_time = datetime.now()

        if not self.enabled:
            return SkillResult(
                success=False,
                error="Perplexity API key not configured",
                skill_name=self.name
            )

        api_key = self._get_api_key("PERPLEXITY")

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }

        model = "llama-3.1-sonar-large-128k-online" if detailed else "llama-3.1-sonar-small-128k-online"

        payload = {
            "model": model,
            "messages": [{"role": "user", "content": query}]
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.perplexity.ai/chat/completions",
                    headers=headers,
                    json=payload
                ) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        response = data["choices"][0]["message"]["content"]
                        citations = data.get("citations", [])

                        execution_time = (datetime.now() - start_time).total_seconds() * 1000

                        result = SkillResult(
                            success=True,
                            data={
                                "response": response,
                                "citations": citations,
                                "model": model
                            },
                            message=response,
                            skill_name=self.name,
                            execution_time_ms=int(execution_time)
                        )

                        self._track_execution(result, execution_time)
                        return result
                    else:
                        error = await resp.text()
                        return SkillResult(
                            success=False,
                            error=f"Perplexity API error: {error}",
                            skill_name=self.name
                        )

        except Exception as e:
            return SkillResult(
                success=False,
                error=str(e),
                skill_name=self.name
            )


# ============================================================================
# PRESENTATION SKILL
# ============================================================================

class PresentationSkill(Skill):
    """Create PowerPoint presentations"""

    name = "presentation"
    description = "Create professional PowerPoint presentations"
    description_it = "Crea presentazioni PowerPoint professionali"
    examples = [
        "Crea una presentazione sul cambiamento climatico",
        "Genera slide per un pitch di startup",
        "Create a presentation about machine learning",
        "Prepare slides for quarterly review"
    ]
    requires_api = []  # Uses local python-pptx
    skill_type = SkillType.GOD
    category = SkillCategory.CREATIVE
    icon = "📊"
    estimated_cost_usd = 0.0

    def __init__(self, api_keys: Optional[Dict[str, str]] = None):
        super().__init__(api_keys)
        try:
            from pptx import Presentation
            self.pptx_available = True
        except ImportError:
            self.pptx_available = False
            self.enabled = False
            logger.warning("python-pptx not installed. Presentation skill disabled.")

    async def execute(
        self,
        title: str,
        slides: List[Dict],
        output_path: Optional[str] = None,
        template: Optional[str] = None
    ) -> SkillResult:
        """
        Create a presentation

        Args:
            title: Presentation title
            slides: List of slide definitions
                   [{"title": "...", "content": ["bullet1", "bullet2"], "notes": "..."}]
            output_path: Where to save (default: temp file)
            template: Optional template path

        Returns:
            SkillResult with file path
        """
        start_time = datetime.now()

        if not self.pptx_available:
            return SkillResult(
                success=False,
                error="python-pptx not installed. Run: pip install python-pptx",
                skill_name=self.name
            )

        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches, Pt

        try:
            prs = PptxPresentation(template) if template else PptxPresentation()

            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = datetime.now().strftime("%Y-%m-%d")

            # Content slides
            bullet_slide_layout = prs.slide_layouts[1]
            for slide_def in slides:
                slide = prs.slides.add_slide(bullet_slide_layout)
                slide.shapes.title.text = slide_def.get("title", "")

                # Add content
                body = slide.shapes.placeholders[1]
                tf = body.text_frame

                content = slide_def.get("content", [])
                for i, bullet in enumerate(content):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0

                # Add notes
                if slide_def.get("notes"):
                    slide.notes_slide.notes_text_frame.text = slide_def["notes"]

            # Save
            if output_path is None:
                output_path = tempfile.mktemp(suffix=".pptx")

            prs.save(output_path)

            execution_time = (datetime.now() - start_time).total_seconds() * 1000

            result = SkillResult(
                success=True,
                data={"path": output_path, "slides": len(slides) + 1},
                message=f"Presentation created with {len(slides) + 1} slides",
                artifacts=[output_path],
                skill_name=self.name,
                execution_time_ms=int(execution_time)
            )

            self._track_execution(result, execution_time)
            return result

        except Exception as e:
            return SkillResult(
                success=False,
                error=str(e),
                skill_name=self.name
            )


# ============================================================================
# DOCUMENT GENERATION SKILL
# ============================================================================

class DocumentGenerationSkill(Skill):
    """Generate Word/PDF documents"""

    name = "document_generation"
    description = "Generate professional documents (Word, PDF)"
    description_it = "Genera documenti professionali (Word, PDF)"
    examples = [
        "Crea un documento con questi contenuti",
        "Genera un report in Word",
        "Create a document from this outline",
        "Generate a professional letter"
    ]
    requires_api = []
    skill_type = SkillType.GOD
    category = SkillCategory.CREATIVE
    icon = "📝"
    estimated_cost_usd = 0.0

    def __init__(self, api_keys: Optional[Dict[str, str]] = None):
        super().__init__(api_keys)
        try:
            from docx import Document
            self.docx_available = True
        except ImportError:
            self.docx_available = False
            self.enabled = False
            logger.warning("python-docx not installed. Document skill disabled.")

    async def execute(
        self,
        title: str,
        sections: List[Dict],
        output_path: Optional[str] = None,
        output_format: str = "docx"
    ) -> SkillResult:
        """
        Create a document

        Args:
            title: Document title
            sections: List of section definitions
                     [{"heading": "...", "content": "..."}]
            output_path: Where to save (default: temp file)
            output_format: "docx" or "pdf"

        Returns:
            SkillResult with file path
        """
        start_time = datetime.now()

        if not self.docx_available:
            return SkillResult(
                success=False,
                error="python-docx not installed. Run: pip install python-docx",
                skill_name=self.name
            )

        from docx import Document

        try:
            doc = Document()

            # Title
            doc.add_heading(title, 0)

            # Sections
            for section in sections:
                if section.get("heading"):
                    doc.add_heading(section["heading"], 1)
                if section.get("content"):
                    doc.add_paragraph(section["content"])

            # Save
            if output_path is None:
                output_path = tempfile.mktemp(suffix=f".{output_format}")

            doc.save(output_path)

            execution_time = (datetime.now() - start_time).total_seconds() * 1000

            result = SkillResult(
                success=True,
                data={"path": output_path, "sections": len(sections)},
                message=f"Document created with {len(sections)} sections",
                artifacts=[output_path],
                skill_name=self.name,
                execution_time_ms=int(execution_time)
            )

            self._track_execution(result, execution_time)
            return result

        except Exception as e:
            return SkillResult(
                success=False,
                error=str(e),
                skill_name=self.name
            )


# ============================================================================
# CODE ANALYSIS SKILL
# ============================================================================

class CodeAnalysisSkill(Skill):
    """Analyze and explain code"""

    name = "code_analysis"
    description = "Analyze code for bugs, security issues, and improvements"
    description_it = "Analizza codice per bug, problemi di sicurezza e miglioramenti"
    examples = [
        "Analizza questo codice Python",
        "Trova bug in questo script",
        "Review this code for security issues",
        "Explain what this function does"
    ]
    requires_api = ["CLAUDE"]  # Uses Claude for analysis
    skill_type = SkillType.GOD
    category = SkillCategory.TECHNICAL
    icon = "💻"
    estimated_cost_usd = 0.01

    async def execute(
        self,
        code: str,
        language: str = "python",
        analysis_type: str = "general"
    ) -> SkillResult:
        """
        Analyze code

        Args:
            code: Source code to analyze
            language: Programming language
            analysis_type: "general", "security", "performance", "bugs"

        Returns:
            SkillResult with analysis
        """
        start_time = datetime.now()

        if not self.enabled:
            return SkillResult(
                success=False,
                error="Claude API key not configured",
                skill_name=self.name
            )

        api_key = self._get_api_key("CLAUDE")

        analysis_prompts = {
            "general": f"Analyze this {language} code. Explain what it does, identify any issues, and suggest improvements.",
            "security": f"Perform a security analysis of this {language} code. Identify potential vulnerabilities and suggest fixes.",
            "performance": f"Analyze this {language} code for performance. Identify bottlenecks and suggest optimizations.",
            "bugs": f"Review this {language} code for bugs. Identify potential issues and edge cases."
        }

        prompt = analysis_prompts.get(analysis_type, analysis_prompts["general"])
        prompt += f"\n\n```{language}\n{code}\n```"

        headers = {
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01",
            "content-type": "application/json"
        }

        payload = {
            "model": "claude-3-5-haiku-20241022",  # Fast model for analysis
            "max_tokens": 2000,
            "messages": [{"role": "user", "content": prompt}]
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.anthropic.com/v1/messages",
                    headers=headers,
                    json=payload
                ) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        analysis = data["content"][0]["text"]

                        execution_time = (datetime.now() - start_time).total_seconds() * 1000

                        result = SkillResult(
                            success=True,
                            data={
                                "analysis": analysis,
                                "language": language,
                                "analysis_type": analysis_type
                            },
                            message=analysis,
                            skill_name=self.name,
                            execution_time_ms=int(execution_time)
                        )

                        self._track_execution(result, execution_time)
                        return result
                    else:
                        error = await resp.text()
                        return SkillResult(
                            success=False,
                            error=f"Claude API error: {error}",
                            skill_name=self.name
                        )

        except Exception as e:
            return SkillResult(
                success=False,
                error=str(e),
                skill_name=self.name
            )


# ============================================================================
# EMAIL DRAFT SKILL
# ============================================================================

class EmailDraftSkill(Skill):
    """Draft professional emails"""

    name = "email_draft"
    description = "Draft professional emails based on context and tone"
    description_it = "Redige email professionali in base al contesto e tono"
    examples = [
        "Scrivi un'email di follow-up",
        "Draft a professional reply to this email",
        "Write a meeting invitation email",
        "Compose a thank you email"
    ]
    requires_api = ["CLAUDE"]
    skill_type = SkillType.GOD
    category = SkillCategory.COMMUNICATION
    icon = "✉️"
    estimated_cost_usd = 0.005

    async def execute(
        self,
        context: str,
        tone: str = "professional",
        recipient: Optional[str] = None,
        subject_hint: Optional[str] = None,
        language: str = "en"
    ) -> SkillResult:
        """
        Draft an email

        Args:
            context: What the email should be about
            tone: "professional", "friendly", "formal", "casual"
            recipient: Who the email is for (optional context)
            subject_hint: Hint for the subject line
            language: "en" or "it"

        Returns:
            SkillResult with email draft
        """
        start_time = datetime.now()

        if not self.enabled:
            return SkillResult(
                success=False,
                error="Claude API key not configured",
                skill_name=self.name
            )

        api_key = self._get_api_key("CLAUDE")

        lang_instruction = "Write in English." if language == "en" else "Scrivi in italiano."

        prompt = f"""Draft a {tone} email based on this context:

Context: {context}
{f"Recipient: {recipient}" if recipient else ""}
{f"Subject should relate to: {subject_hint}" if subject_hint else ""}

{lang_instruction}

Provide the email in this format:
Subject: [subject line]

[email body]"""

        headers = {
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01",
            "content-type": "application/json"
        }

        payload = {
            "model": "claude-3-5-haiku-20241022",
            "max_tokens": 1000,
            "messages": [{"role": "user", "content": prompt}]
        }

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    "https://api.anthropic.com/v1/messages",
                    headers=headers,
                    json=payload
                ) as resp:
                    if resp.status == 200:
                        data = await resp.json()
                        draft = data["content"][0]["text"]

                        # Parse subject and body
                        lines = draft.strip().split("\n")
                        subject = ""
                        body = draft
                        if lines[0].lower().startswith("subject:"):
                            subject = lines[0].replace("Subject:", "").strip()
                            body = "\n".join(lines[2:]).strip()

                        execution_time = (datetime.now() - start_time).total_seconds() * 1000

                        result = SkillResult(
                            success=True,
                            data={
                                "subject": subject,
                                "body": body,
                                "full_draft": draft,
                                "tone": tone
                            },
                            message=draft,
                            skill_name=self.name,
                            execution_time_ms=int(execution_time)
                        )

                        self._track_execution(result, execution_time)
                        return result
                    else:
                        error = await resp.text()
                        return SkillResult(
                            success=False,
                            error=f"Claude API error: {error}",
                            skill_name=self.name
                        )

        except Exception as e:
            return SkillResult(
                success=False,
                error=str(e),
                skill_name=self.name
            )


# ============================================================================
# CALENDAR SKILL
# ============================================================================

class CalendarSkill(Skill):
    """Manage calendar events"""

    name = "calendar"
    description = "Create and manage calendar events"
    description_it = "Crea e gestisci eventi del calendario"
    examples = [
        "Crea un evento per domani alle 10",
        "Schedule a meeting for next Monday",
        "Add a reminder for the project deadline",
        "Create a recurring weekly meeting"
    ]
    requires_api = []  # Integration with OAuth calendar providers
    skill_type = SkillType.GOD
    category = SkillCategory.COMMUNICATION
    icon = "📅"
    estimated_cost_usd = 0.0

    async def execute(
        self,
        action: str,
        title: Optional[str] = None,
        start_time: Optional[str] = None,
        end_time: Optional[str] = None,
        description: Optional[str] = None,
        attendees: Optional[List[str]] = None,
        **kwargs
    ) -> SkillResult:
        """
        Manage calendar events

        Args:
            action: "create", "list", "update", "delete"
            title: Event title
            start_time: ISO format start time
            end_time: ISO format end time
            description: Event description
            attendees: List of attendee emails

        Returns:
            SkillResult with calendar operation result
        """
        start = datetime.now()

        # This is a placeholder - actual implementation requires
        # OAuth tokens from the user's connected calendar
        result = SkillResult(
            success=False,
            error="Calendar integration requires connected Google/Outlook account",
            skill_name=self.name,
            data={
                "action": action,
                "title": title,
                "start_time": start_time,
                "note": "Connect your calendar in Settings to enable this feature"
            }
        )

        execution_time = (datetime.now() - start).total_seconds() * 1000
        result.execution_time_ms = int(execution_time)

        return result


# ============================================================================
# SKILL REGISTRY
# ============================================================================

GOD_SKILLS = [
    ImageGenerationSkill,
    WebSearchSkill,
    PresentationSkill,
    DocumentGenerationSkill,
    CodeAnalysisSkill,
    EmailDraftSkill,
    CalendarSkill,
]


def get_all_god_skills(api_keys: Optional[Dict[str, str]] = None) -> List[Skill]:
    """Get instances of all GOD skills"""
    return [skill_class(api_keys) for skill_class in GOD_SKILLS]
